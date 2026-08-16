# -*- coding: utf-8 -*-
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path(r"C:\Users\shuai\Documents\我的坚果云\勘探院材料\2026-08-岗位实践大赛\岗位实践大赛-郑帅-2.pptx")
OUTPUT = SOURCE.with_name("岗位实践大赛-郑帅-2-补充完善版.pptx")
TEAM_SCREENSHOT = Path(r"C:\Users\shuai\Documents\QwenPaw\plugins\bundle\ugsci\docs\assets\screenshots\07-expert-teams.png")

RED = RGBColor(198, 0, 0)
DARK_RED = RGBColor(146, 20, 20)
BLUE = RGBColor(57, 113, 178)
LIGHT_BLUE = RGBColor(225, 238, 250)
GREEN = RGBColor(76, 145, 65)
LIGHT_GREEN = RGBColor(229, 242, 223)
ORANGE = RGBColor(235, 126, 36)
LIGHT_ORANGE = RGBColor(253, 235, 220)
GOLD = RGBColor(239, 178, 25)
LIGHT_GOLD = RGBColor(255, 244, 205)
GRAY = RGBColor(92, 92, 92)
LIGHT_GRAY = RGBColor(244, 245, 247)
BLACK = RGBColor(24, 24, 24)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def set_text(
    shape, text: str, size: float, color=BLACK, bold=False,
    align=PP_ALIGN.LEFT, font="Microsoft YaHei",
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color


def text_box(
    slide, x, y, w, h, text, size=18, color=BLACK, bold=False,
    align=PP_ALIGN.LEFT, fill=None, line=None, radius=True,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill if fill else WHITE
    shp.line.color.rgb = line if line else (fill if fill else WHITE)
    shp.line.width = Pt(1)
    set_text(shp, text, size, color, bold, align)
    return shp


def add_header(slide, title: str, section="● 智能体搭建") -> None:
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(12.75), Inches(0.64))
    set_text(title_box, title, 30, RED, True, PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.82), Inches(13.33), Inches(0.035))
    line.fill.solid(); line.fill.fore_color.rgb = RED; line.line.color.rgb = RED
    section_box = slide.shapes.add_textbox(Inches(0.42), Inches(0.93), Inches(4.5), Inches(0.43))
    set_text(section_box, section, 21, RED, True)


def add_footer(slide, text="数据来源：呼图壁储气库2024—2025周期自评材料；计算结果为建议复核值") -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(7.44), Inches(13.33), Inches(0.025))
    line.fill.solid(); line.fill.fore_color.rgb = DARK_RED; line.line.color.rgb = DARK_RED
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.20), Inches(12.2), Inches(0.2))
    set_text(box, text, 8.5, GRAY, False, PP_ALIGN.RIGHT)


def arrow(slide, x, y, w=0.38, h=0.34, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.color.rgb = color
    return shp


def build_slide_9(slide) -> None:
    clear_slide(slide)
    add_header(slide, "智能体搭建：从业务流程到可执行智能体")
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.37), Inches(12.0), Inches(0.42))
    set_text(sub, "将行业标准、专家经验、确定性算法和质量门控封装成可复用的库存评价智能体", 17, BLACK, True)

    cards = [
        ("① 任务拆解", "明确数据输入、状态边界\n拆分库存、工作气、冲峰指标", ORANGE, LIGHT_ORANGE),
        ("② 能力装配", "配置专家角色与业务方法\n挂接领域知识和算法模型", GOLD, LIGHT_GOLD),
        ("③ 流程编排", "领导Agent理解任务\n多专家协同、自动选择方法", BLUE, LIGHT_BLUE),
        ("④ 质量门控", "单位与压力口径校验\n结果状态、来源与指纹可追溯", GREEN, LIGHT_GREEN),
    ]
    x_positions = [0.65, 3.75, 6.85, 9.95]
    for i, ((title, body, color, light), x) in enumerate(zip(cards, x_positions)):
        text_box(slide, x, 1.95, 2.72, 0.48, title, 16, WHITE, True, PP_ALIGN.CENTER, color, color)
        text_box(slide, x, 2.42, 2.72, 1.18, body, 13.5, BLACK, False, PP_ALIGN.CENTER, light, color)
        if i < 3:
            arrow(slide, x + 2.78, 2.74, 0.28, 0.28, RED)

    slide.shapes.add_picture(str(TEAM_SCREENSHOT), Inches(0.72), Inches(3.88), width=Inches(6.15), height=Inches(3.08))
    text_box(slide, 6.98, 3.88, 5.65, 0.45, "库存评价智能体的最小定义", 17, WHITE, True, PP_ALIGN.CENTER, RED, RED)
    items = [
        ("做什么", "读取资料，识别评价任务与周期边界"),
        ("怎么做", "分层计算、账面核算、指标对照、交叉验证"),
        ("如何协同", "资料审核 → 动态分析 → 库存评价 → 结果复核"),
        ("如何保证", "确定性计算负责数值；Agent负责理解、编排与解释"),
    ]
    y = 4.46
    colors = [ORANGE, GOLD, BLUE, GREEN]
    lights = [LIGHT_ORANGE, LIGHT_GOLD, LIGHT_BLUE, LIGHT_GREEN]
    for (label, desc), color, light in zip(items, colors, lights):
        text_box(slide, 7.02, y, 1.15, 0.50, label, 13.5, WHITE, True, PP_ALIGN.CENTER, color, color)
        text_box(slide, 8.20, y, 4.40, 0.50, desc, 12.5, BLACK, False, PP_ALIGN.LEFT, light, color)
        y += 0.62
    add_footer(slide, "搭建原则：专家定义保持简洁，不绑定具体工具名称；平台按任务自动路由能力")


def build_slide_10(slide) -> None:
    clear_slide(slide)
    add_header(slide, "场景应用：呼图壁24—25周期运行实例")
    text_box(slide, 0.62, 1.42, 3.05, 0.43, "任务输入", 17, WHITE, True, PP_ALIGN.CENTER, ORANGE, ORANGE)
    input_text = (
        "评价对象：呼图壁储气库\n"
        "材料：2024—2025周期自评PPT\n"
        "设计库容：107.0 亿方\n"
        "账面库存：104.88 亿方\n"
        "工作气量：43.5 / 45.1 亿方\n"
        "冲峰能力：3950 / 4020 万方/日\n"
        "压力口径：视地层压力（显式声明）"
    )
    text_box(slide, 0.62, 1.86, 3.05, 3.55, input_text, 12.2, BLACK, False, PP_ALIGN.LEFT, LIGHT_ORANGE, ORANGE)
    text_box(slide, 0.62, 5.60, 3.05, 1.05, "用户只需描述任务并上传材料\n无需指定具体工具名称", 15, RED, True, PP_ALIGN.CENTER, WHITE, RED)

    text_box(slide, 3.98, 1.42, 5.34, 0.43, "智能体自动执行链", 17, WHITE, True, PP_ALIGN.CENTER, BLUE, BLUE)
    stages = [
        ("① 资料抽取", "识别层系、周期、气量、压力与Z因子"),
        ("② 口径校验", "统一单位、时间边界和视地层压力口径"),
        ("③ 专家协同", "资料审核、动态分析、库存评价；领导Agent复核"),
        ("④ 确定性计算", "分层p/Z计算后汇总；账面库存独立核算"),
        ("⑤ 评价输出", "指标对照、风险提示、报告和全过程溯源"),
    ]
    y = 1.98
    colors = [ORANGE, GOLD, BLUE, GREEN, RED]
    lights = [LIGHT_ORANGE, LIGHT_GOLD, LIGHT_BLUE, LIGHT_GREEN, RGBColor(252, 229, 229)]
    for idx, ((name, desc), color, light) in enumerate(zip(stages, colors, lights)):
        text_box(slide, 4.08, y, 1.42, 0.58, name, 13, WHITE, True, PP_ALIGN.CENTER, color, color)
        text_box(slide, 5.52, y, 3.68, 0.58, desc, 12.2, BLACK, False, PP_ALIGN.LEFT, light, color)
        if idx < len(stages) - 1:
            down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.38), Inches(y + 0.58), Inches(0.32), Inches(0.22))
            down.fill.solid(); down.fill.fore_color.rgb = GRAY; down.line.color.rgb = GRAY
        y += 0.91

    text_box(slide, 9.62, 1.42, 3.05, 0.43, "自动调用与门控", 17, WHITE, True, PP_ALIGN.CENTER, GREEN, GREEN)
    calls = [
        ("3类专家", "按专业分工协同，不重复手算"),
        ("1个综合入口", "一次返回分层结果和评价指标"),
        ("5类校验", "单位、状态、口径、公式、结果状态"),
        ("全程留痕", "输入指纹、公式、假设、告警可追溯"),
    ]
    y = 1.95
    for label, desc in calls:
        text_box(slide, 9.78, y, 1.08, 0.75, label, 15, RED, True, PP_ALIGN.CENTER, WHITE, RED)
        text_box(slide, 10.92, y, 1.56, 0.75, desc, 11.3, BLACK, False, PP_ALIGN.CENTER, LIGHT_GREEN, GREEN)
        y += 0.92
    text_box(slide, 9.78, 5.80, 2.70, 0.88, "输出固定标记：\n计算建议值，待业务复核", 15, RED, True, PP_ALIGN.CENTER, RGBColor(255, 244, 244), RED)
    add_footer(slide, "说明：分层采气量与Z因子按报告发布结果重构，用于演示确定性回归验算；并非替代原始测试数据")


def build_slide_11(slide) -> None:
    clear_slide(slide)
    add_header(slide, "场景应用：呼图壁结果对比与确定性验证")

    text_box(slide, 0.62, 1.42, 4.05, 0.43, "分层有效库存计算结果", 17, WHITE, True, PP_ALIGN.CENTER, BLUE, BLUE)
    kpis = [
        ("E1-2z21", "78.3", "亿方"),
        ("E1-2z22", "26.7", "亿方"),
        ("合计建议复核值", "105.0", "亿方"),
    ]
    y = 1.98
    for i, (label, value, unit) in enumerate(kpis):
        color = [ORANGE, GOLD, RED][i]
        light = [LIGHT_ORANGE, LIGHT_GOLD, RGBColor(255, 235, 235)][i]
        text_box(slide, 0.72, y, 1.65, 0.82, label, 13, WHITE, True, PP_ALIGN.CENTER, color, color)
        text_box(slide, 2.40, y, 1.48, 0.82, value, 24, color, True, PP_ALIGN.CENTER, light, color)
        text_box(slide, 3.91, y, 0.60, 0.82, unit, 12, BLACK, True, PP_ALIGN.CENTER, WHITE, color)
        y += 1.03
    text_box(slide, 0.72, 5.24, 3.79, 1.14, "计算口径\n逐层计算后汇总；压力基准为视地层压力\n不使用全库平均压力替代分层计算", 12.7, BLACK, False, PP_ALIGN.CENTER, LIGHT_BLUE, BLUE)

    text_box(slide, 4.92, 1.42, 5.12, 0.43, "材料结果与智能体结果对比", 17, WHITE, True, PP_ALIGN.CENTER, GREEN, GREEN)
    headers = ["指标", "材料/人工", "智能体", "差异"]
    rows = [
        ["E1-2z21", "78.3", "78.3", "0.0"],
        ["E1-2z22", "26.7", "26.7", "0.0"],
        ["有效库存合计", "105.0", "105.0", "0.0"],
        ["工作气符合率", "约96.5%", "96.45%", "<0.05%"],
        ["冲峰能力符合率", "约98.3%", "98.26%", "<0.05%"],
    ]
    col_x = [4.98, 6.65, 7.89, 9.08]
    col_w = [1.62, 1.19, 1.14, 0.88]
    for x, w, h in zip(col_x, col_w, headers):
        text_box(slide, x, 1.93, w, 0.48, h, 12.5, WHITE, True, PP_ALIGN.CENTER, GREEN, GREEN)
    y = 2.42
    for ridx, row in enumerate(rows):
        fill = WHITE if ridx % 2 == 0 else LIGHT_GREEN
        for x, w, cell in zip(col_x, col_w, row):
            text_box(slide, x, y, w, 0.55, cell, 11.8, BLACK, ridx == 2, PP_ALIGN.CENTER, fill, RGBColor(183, 205, 178))
        y += 0.55
    text_box(slide, 5.02, 5.39, 4.92, 0.99, "验证结论：核心库存结果与材料一致\n计算过程可重复，公式、输入和告警均可追踪", 14.5, GREEN, True, PP_ALIGN.CENTER, LIGHT_GREEN, GREEN)

    text_box(slide, 10.30, 1.42, 2.37, 0.43, "确定性保证", 17, WHITE, True, PP_ALIGN.CENTER, RED, RED)
    guarantees = [
        ("公式固定", "数值计算不由大模型自由生成"),
        ("边界固定", "层系、周期、单位和压力口径显式校验"),
        ("状态固定", "105.0亿方只标记为建议复核值"),
        ("证据固定", "保留来源、输入指纹、公式与告警"),
    ]
    y = 1.98
    for label, desc in guarantees:
        text_box(slide, 10.43, y, 0.85, 0.72, label, 12, WHITE, True, PP_ALIGN.CENTER, RED, RED)
        text_box(slide, 11.31, y, 1.23, 0.72, desc, 10.8, BLACK, False, PP_ALIGN.CENTER, RGBColor(255, 241, 241), RED)
        y += 0.89
    text_box(slide, 10.43, 5.75, 2.11, 0.63, "账面达容率 98.02%", 14, RED, True, PP_ALIGN.CENTER, WHITE, RED)
    add_footer(slide, "报告发布值：E1-2z21=78.3、E1-2z22=26.7、合计105.0亿方；演示输入用于公式回归验证")


def build_slide_13(slide) -> None:
    clear_slide(slide)
    add_header(slide, "成果展示：UGSci平台与可复用能力资产", "● 储气库智能体平台（UGSci）")
    slide.shapes.add_picture(str(TEAM_SCREENSHOT), Inches(0.63), Inches(1.50), width=Inches(7.38), height=Inches(4.16))
    text_box(slide, 0.63, 5.82, 7.38, 0.60, "专家、专家团和工作流可视化配置，任务自动分派并保留调用过程", 14, BLUE, True, PP_ALIGN.CENTER, LIGHT_BLUE, BLUE)

    text_box(slide, 8.30, 1.50, 4.38, 0.43, "形成的可复用资产", 17, WHITE, True, PP_ALIGN.CENTER, RED, RED)
    assets = [
        ("知识资产", "行业规范、报告口径、专家经验"),
        ("算法资产", "物质平衡、库存指标、运行指标"),
        ("智能体资产", "任务理解、专业协同、结果复核"),
        ("流程资产", "取数—计算—评价—报告全链路"),
    ]
    colors = [ORANGE, GOLD, BLUE, GREEN]
    lights = [LIGHT_ORANGE, LIGHT_GOLD, LIGHT_BLUE, LIGHT_GREEN]
    y = 2.05
    for (label, desc), color, light in zip(assets, colors, lights):
        text_box(slide, 8.42, y, 1.08, 0.72, label, 12.5, WHITE, True, PP_ALIGN.CENTER, color, color)
        text_box(slide, 9.54, y, 2.98, 0.72, desc, 12, BLACK, False, PP_ALIGN.CENTER, light, color)
        y += 0.88
    text_box(slide, 8.42, 5.69, 4.10, 0.76, "平台优势：本地运行、数据不出域\n专家定义简洁，领域能力持续复用", 14, RED, True, PP_ALIGN.CENTER, RGBColor(255, 243, 243), RED)
    add_footer(slide, "成果不止是一次计算，而是一套可持续沉淀、复用和扩展的领域智能体平台")


def build_slide_14(slide) -> None:
    clear_slide(slide)
    add_header(slide, "成果感悟：数智建功，岗位成才")
    cards = [
        ("业务价值", ["把分散核算转化为标准化流程", "核心结果与材料一致、全过程可追溯", "为调峰保供提供及时、可解释的评价依据"], RED, RGBColor(255, 239, 239)),
        ("方法创新", ["大模型负责理解和编排", "确定性工具负责可靠计算", "多专家交叉验证，关键口径显式门控"], BLUE, LIGHT_BLUE),
        ("岗位成长", ["从单项技术研究走向系统解决方案", "把个人经验沉淀为团队可复用资产", "以数字化能力服务能源保供与安全生产"], GREEN, LIGHT_GREEN),
    ]
    x_positions = [0.72, 4.56, 8.40]
    for (title, bullets, color, light), x in zip(cards, x_positions):
        text_box(slide, x, 1.48, 3.32, 0.52, title, 18, WHITE, True, PP_ALIGN.CENTER, color, color)
        body = "\n".join(f"● {b}" for b in bullets)
        text_box(slide, x, 2.02, 3.32, 2.40, body, 14, BLACK, False, PP_ALIGN.LEFT, light, color)

    text_box(slide, 0.72, 4.76, 11.00, 0.48, "下一步：从“可用算例”走向“规模化应用”", 18, WHITE, True, PP_ALIGN.CENTER, DARK_RED, DARK_RED)
    next_steps = [
        ("补齐数据入口", "对接生产数据库与监测数据，减少人工取数"),
        ("扩大算例覆盖", "建立多库、多周期黄金数据集和回归测试"),
        ("完善评价闭环", "接入业务复核、批准和报告归档流程"),
    ]
    y = 5.43
    x = 0.72
    for idx, (label, desc) in enumerate(next_steps, 1):
        color = [ORANGE, BLUE, GREEN][idx - 1]
        light = [LIGHT_ORANGE, LIGHT_BLUE, LIGHT_GREEN][idx - 1]
        text_box(slide, x, y, 0.47, 0.80, str(idx), 20, WHITE, True, PP_ALIGN.CENTER, color, color)
        text_box(slide, x + 0.50, y, 1.35, 0.80, label, 13, color, True, PP_ALIGN.CENTER, WHITE, color)
        text_box(slide, x + 1.88, y, 1.55, 0.80, desc, 11.2, BLACK, False, PP_ALIGN.CENTER, light, color)
        x += 3.68
    add_footer(slide, "建议赛前再补充一组真实耗时数据：人工评价耗时 vs 智能体运行耗时")


def main() -> None:
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(OUTPUT)
    build_slide_9(prs.slides[8])
    build_slide_10(prs.slides[9])
    build_slide_11(prs.slides[10])
    build_slide_13(prs.slides[12])
    build_slide_14(prs.slides[13])
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
