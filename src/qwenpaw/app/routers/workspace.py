# -*- coding: utf-8 -*-
"""Workspace API – download / upload the entire WORKING_DIR as a zip.

Also includes agent file management, language settings, audio/transcription
configuration, running config, and system prompt files.
"""

from __future__ import annotations

import asyncio
import copy
import io
import json
import logging
import mimetypes
import secrets
import shutil
import stat
import os
import subprocess
import sys
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, AsyncIterator, Literal
from urllib.parse import parse_qs, quote, unquote, urlparse

from fastapi import (
    APIRouter,
    Body,
    File,
    HTTPException,
    Query,
    Request,
    UploadFile,
)
from fastapi.responses import (
    FileResponse,
    ORJSONResponse,
    Response,
    StreamingResponse,
)
from watchfiles import awatch, Change
from pydantic import BaseModel, Field

from ..utils import check_upload_size, safe_join, schedule_agent_reload
from ...config import (
    load_config,
    AgentsRunningConfig,
)
from ...config.utils import mutate_config
from ...config.config import (
    EmbeddingModelConfig,
    load_agent_config,
    save_agent_config,
    update_agent_config_async,
)
from ...agents.memory.embedding_model import (
    embedding_vector_space_fingerprint,
    test_embedding_model,
)
from ...agents.memory.agent_md_manager import AgentMdManager
from ...agents.templates import get_workspace_md_template_id
from ...agents.utils import copy_workspace_md_files
from ...constant import (
    BUILTIN_QA_AGENT_ID,
    SUPPORTED_AGENT_LANGUAGES,
    WORKING_DIR,
)
from ...services.fs_name_rules import NameRules, probe_name_rules
from ...services.reveal_file import reveal_workspace_path
from ...services.workspace_files import (
    DEFAULT_CHUNK_SIZE,
    DEFAULT_PAGE_SIZE,
    FileVersionConflict,
    InvalidCursor,
    InvalidWorkspacePath,
    MAX_PAGE_SIZE,
    file_etag,
    get_file_metadata,
    list_directory,
    read_file_chunk,
    resolve_workspace_path,
    save_text_file,
)
from ...utils.io_utils import get_path_lock, run_sync_io
from ..agent_context import (
    get_agent_for_request,
    get_agent_project_dir,
    get_project_dir_for_request,
    get_project_dirs_for_request,
)

router = APIRouter(prefix="/workspace", tags=["workspace"])
logger = logging.getLogger(__name__)
_FILESYSTEM_SEMAPHORE = asyncio.Semaphore(8)
_WATCH_HEARTBEAT_SECONDS = 30.0
_WATCH_POLL_TIMEOUT_MS = 1_000


class MdFileInfo(BaseModel):
    """Markdown file metadata."""

    filename: str = Field(..., description="File name")
    path: str = Field(..., description="File path")
    size: int = Field(..., description="Size in bytes")
    created_time: str = Field(..., description="Created time")
    modified_time: str = Field(..., description="Modified time")


class MdFileContent(BaseModel):
    """Markdown file content."""

    content: str = Field(..., description="File content")


class PromptFileWriteRequest(BaseModel):
    """Create/update a root Markdown file and optionally change
    prompt mount."""

    filename: str = Field(..., description="Portable root Markdown filename")
    content: str = Field(..., description="Markdown content")
    enable: bool | None = Field(
        None,
        description="True to mount, false to unmount, null to preserve config",
    )


class EmbeddingTestResponse(BaseModel):
    """Result of an AgentScope embedding connectivity request."""

    success: bool
    configured_dimensions: int
    actual_dimensions: int | None = None
    latency_ms: int
    message: str


def _dir_stats(root: Path) -> tuple[int, int]:
    """Return (file_count, total_size) for *root* recursively."""
    count = 0
    size = 0
    if root.is_dir():
        for p in root.rglob("*"):
            if p.is_file():
                count += 1
                size += p.stat().st_size
    return count, size


def _zip_directory(root: Path) -> io.BytesIO:
    """Create an in-memory zip archive of *root* and return the buffer.

    All files **and** directories (including empty ones) are included.
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for entry in sorted(root.rglob("*")):
            arcname = entry.relative_to(root).as_posix()
            if entry.is_file():
                zf.write(entry, arcname)
            elif entry.is_dir():
                # Zip spec: directory entries end with '/'
                zf.write(entry, arcname + "/")
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Agent File Management Endpoints
# ---------------------------------------------------------------------------


@router.get(
    "/files",
    response_model=list[MdFileInfo],
    summary="List working files",
    description="List all working files (uses active agent)",
)
async def list_working_files(
    request: Request,
) -> list[MdFileInfo]:
    """List working directory markdown files."""
    try:
        workspace = await get_agent_for_request(request)
        workspace_manager = AgentMdManager(
            str(workspace.workspace_dir),
            agent_id=workspace.agent_id,
        )
        files = [
            MdFileInfo.model_validate(file)
            for file in workspace_manager.list_working_mds()
        ]
        return files
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.get(
    "/files/{md_name}",
    response_model=MdFileContent,
    summary="Read a working file",
    description="Read a working markdown file (uses active agent)",
)
async def read_working_file(
    md_name: str,
    request: Request,
) -> MdFileContent:
    """Read a working directory markdown file."""
    try:
        workspace = await get_agent_for_request(request)
        workspace_manager = AgentMdManager(
            str(workspace.workspace_dir),
            agent_id=workspace.agent_id,
        )
        content = workspace_manager.read_working_md(md_name)
        return MdFileContent(content=content)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.put(
    "/files/{md_name}",
    response_model=dict,
    summary="Write a working file",
    description="Create or update a working file (uses active agent)",
)
async def write_working_file(
    md_name: str,
    body: MdFileContent,
    request: Request,
) -> dict:
    """Write a working directory markdown file."""
    try:
        workspace = await get_agent_for_request(request)
        workspace_manager = AgentMdManager(
            str(workspace.workspace_dir),
            agent_id=workspace.agent_id,
        )
        workspace_manager.write_working_md(md_name, body.content)
        return {"written": True}
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.post(
    "/prompt-files",
    response_model=dict,
    summary="Write a Markdown file and atomically update prompt mounting",
)
async def write_prompt_file(
    body: PromptFileWriteRequest,
    request: Request,
) -> dict:
    """Write Markdown and compensate the file if config persistence fails."""
    workspace = await get_agent_for_request(request)
    manager = AgentMdManager(
        str(workspace.workspace_dir),
        agent_id=workspace.agent_id,
    )
    try:
        filename = manager.normalize_working_md_name(body.filename)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    file_path = manager.working_dir / filename
    # pylint: disable=protected-access
    manager._assert_within_dir(file_path, manager.working_dir)
    existed = file_path.exists()
    previous_bytes = file_path.read_bytes() if existed else None
    agent_config = load_agent_config(workspace.agent_id)
    previous_prompt_files = list(agent_config.system_prompt_files or [])
    next_prompt_files = list(previous_prompt_files)

    if body.enable is True and filename not in next_prompt_files:
        next_prompt_files.append(filename)
    elif body.enable is False:
        next_prompt_files = [
            name for name in next_prompt_files if name != filename
        ]

    try:
        manager.write_working_md(filename, body.content)
        if body.enable is not None:
            agent_config.system_prompt_files = next_prompt_files
            save_agent_config(workspace.agent_id, agent_config)
    except Exception as exc:
        try:
            if existed and previous_bytes is not None:
                file_path.write_bytes(previous_bytes)
            elif file_path.exists():
                file_path.unlink()
        except OSError:
            logger.exception("Failed to roll back prompt file %s", filename)
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    try:
        schedule_agent_reload(request, workspace.agent_id)
    except Exception:
        logger.exception("Failed to schedule reload after saving %s", filename)

    return {
        "written": True,
        "filename": filename,
        "system_prompt_files": next_prompt_files,
    }


# ---------------------------------------------------------------------------
# Coding Mode – full file-tree + file watcher (SSE)
# ---------------------------------------------------------------------------

_SKIP_NAMES: frozenset[str] = frozenset(
    {
        ".git",
        "__pycache__",
        ".venv",
        "node_modules",
        ".mypy_cache",
        ".pytest_cache",
        ".ruff_cache",
        ".hypothesis",
    },
)


def _should_skip(rel_parts: tuple[str, ...]) -> bool:
    return any(p.startswith(".") or p in _SKIP_NAMES for p in rel_parts)


def _is_skipped_name(name: str) -> bool:
    return name.startswith(".") or name in _SKIP_NAMES


def _list_all_files(workspace_dir: Path) -> list[dict]:
    """Recursively list all non-hidden workspace files.

    Uses ``os.walk(topdown=True)`` and prunes ``dirnames`` in place so that
    we never descend into ``node_modules`` / ``.venv`` / ``.git`` etc. — the
    previous ``Path.rglob('*')`` walked them fully and filtered after the
    fact, which is the dominant cost on real projects. Each file is stat'd
    exactly once. Paths are returned with POSIX ``/`` separators so the
    frontend ``buildTree`` (which splits on ``/``) works on Windows too.
    """
    files: list[dict] = []
    root = str(workspace_dir)
    try:
        for dirpath, dirnames, filenames in os.walk(root, topdown=True):
            # Prune in place — must mutate, not rebind, for os.walk to honor.
            dirnames[:] = sorted(
                d for d in dirnames if not _is_skipped_name(d)
            )
            rel_dir = os.path.relpath(dirpath, root)
            for name in sorted(filenames):
                if _is_skipped_name(name):
                    continue
                full = os.path.join(dirpath, name)
                try:
                    st = os.stat(full)
                except OSError:
                    continue
                rel = (
                    name
                    if rel_dir == "."
                    else f"{rel_dir}/{name}".replace(os.sep, "/")
                )
                files.append(
                    {
                        "filename": rel,
                        "path": rel,
                        "size": st.st_size,
                        "modified_time": datetime.fromtimestamp(
                            st.st_mtime,
                            tz=timezone.utc,
                        ).isoformat(),
                    },
                )
    except Exception:
        pass
    return files


# Prefix selecting a non-primary bound project directory by absolute path,
# e.g. ``project:/Users/me/docs``. The path is carried rather than an index
# because the bound list is reorderable ("make primary"): an index would let a
# persisted editor tab silently start pointing at a different directory.
_EXTRA_PROJECT_ROOT_PREFIX = "project:"


async def _resolve_extra_project_root(
    request: Request,
    workspace: Any,
    raw_path: str,
) -> Path:
    """Resolve one bound project directory selected by absolute path.

    The membership check is the authorization boundary for the Files API: a
    path is served only when it is one of the directories this chat actually
    bound. Anything else is rejected outright — never silently downgraded to
    the primary, which would make an out-of-bounds request look like it
    succeeded against the wrong directory.

    Membership is decided by :func:`dir_key`: the candidate is keyed in a
    worker thread and the loop compares strings, touching the filesystem
    not at all. Two things depend on that split.

    The loop must do no I/O. The obvious spelling —
    ``same_dir(candidate, entry.path)`` — resolves *both* sides on every
    iteration, so ten bound directories cost twenty ``resolve()`` calls on
    the event loop per Files request, half of them re-resolving
    ``entry.path``, which ``ResolvedProjectDirs`` already canonicalized.
    One stalled SMB or FUSE mount in the list would then stall every other
    request the process is serving.

    And the comparison must be by directory identity, not by path text. A
    string comparison decides membership on spelling: fold case and an
    unbound ``/srv/REPO`` is served as ``/srv/repo`` on a case-sensitive
    volume; do not fold and a bound directory reached by a symlink, a
    mount alias or a ``..`` detour is refused with 403. Identity is right
    in both directions without knowing anything about the volume.

    A configured directory that does not exist has no identity, so its key
    is its path text and the comparison degrades to the old spelling-based
    one — acceptable, because there is nothing there to serve either way.
    """
    from ...services.project_directory import dir_key

    candidate = raw_path.strip()
    if not candidate:
        raise HTTPException(status_code=400, detail="root path is empty")

    resolved = await get_project_dirs_for_request(request, workspace)
    candidate_key = await run_sync_io(dir_key, candidate)
    for entry in resolved.dirs:
        # An entry built without a key would otherwise match the empty
        # string; only a real key can grant membership.
        if entry.key and entry.key == candidate_key:
            return entry.path
    # The workspace is a legitimate root, but it has its own ``root=workspace``
    # selector; accepting it here too would let one root be addressed two ways.
    raise HTTPException(
        status_code=403,
        detail="Not a bound project directory",
    )


async def _resolve_files_root(
    request: Request,
    workspace: Any,
    root: str,
) -> Path:
    """Resolve the selected project or agent configuration directory.

    Accepted values:

    * ``workspace`` — the agent's own storage root
    * ``project`` — the PRIMARY bound project directory
    * ``project:<absolute path>`` — any other directory bound to this chat
    """
    if root == "workspace":
        return workspace.workspace_dir
    if root == "project":
        return await get_project_dir_for_request(request, workspace)
    if root.startswith(_EXTRA_PROJECT_ROOT_PREFIX):
        return await _resolve_extra_project_root(
            request,
            workspace,
            root[len(_EXTRA_PROJECT_ROOT_PREFIX) :],
        )
    raise HTTPException(
        status_code=400,
        detail="root must be project, project:<path> or workspace",
    )


@router.get(
    "/tree",
    summary="List one workspace directory page",
)
async def list_workspace_tree(
    request: Request,
    path: str = Query(default=""),
    cursor: str | None = Query(default=None),
    root: str = Query(default="project"),
    limit: int = Query(
        default=DEFAULT_PAGE_SIZE,
        ge=1,
        le=MAX_PAGE_SIZE,
    ),
) -> dict:
    """List immediate children without materializing the full project."""
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)
    try:
        async with _FILESYSTEM_SEMAPHORE:
            return await asyncio.to_thread(
                list_directory,
                files_root,
                path,
                cursor,
                limit,
            )
    except InvalidCursor as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except (FileNotFoundError, NotADirectoryError) as exc:
        raise HTTPException(
            status_code=404,
            detail="Directory not found",
        ) from exc


@router.get(
    "/file-metadata",
    summary="Read workspace file metadata",
)
async def read_workspace_file_metadata(
    request: Request,
    path: str = Query(...),
    root: str = Query(default="project"),
) -> dict:
    """Return file metadata before content is requested."""
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)
    try:
        async with _FILESYSTEM_SEMAPHORE:
            return await asyncio.to_thread(
                get_file_metadata,
                files_root,
                path,
            )
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except (FileNotFoundError, OSError) as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc


def _reveal_allowed_roots(workspace: Any, project_dir: Path) -> list[Path]:
    """Deduplicate project, agent workspace, and global working directories."""
    roots: list[Path] = []
    seen: set[str] = set()
    for candidate in (
        project_dir,
        getattr(workspace, "workspace_dir", None),
        WORKING_DIR,
    ):
        if candidate is None:
            continue
        try:
            resolved = str(Path(candidate).resolve())
        except OSError:
            continue
        if resolved in seen:
            continue
        seen.add(resolved)
        roots.append(Path(resolved))
    return roots


@router.post(
    "/reveal",
    summary="Reveal a workspace file in the OS file manager",
)
async def reveal_workspace_file(
    request: Request,
    path: str = Query(...),
    root: str = Query(default="project"),
) -> dict[str, bool]:
    """Open Explorer / Finder for a project or workspace file.

    Relative POSIX paths resolve under the selected *root*. Absolute host
    paths and ``file:`` URIs are accepted only inside the session project,
    agent workspace, or global working directory.
    """
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)
    project_dir = await get_project_dir_for_request(request, workspace)
    extra_roots = _reveal_allowed_roots(workspace, project_dir)
    try:
        async with _FILESYSTEM_SEMAPHORE:
            await asyncio.to_thread(
                reveal_workspace_path,
                files_root,
                path,
                extra_roots,
            )
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc
    except OSError as exc:
        raise HTTPException(
            status_code=500,
            detail="Unable to open the file manager",
        ) from exc
    return {"ok": True}


@router.get(
    "/file-content",
    summary="Read a bounded workspace text chunk",
)
async def read_workspace_file_content(
    request: Request,
    path: str = Query(...),
    root: str = Query(default="project"),
    offset: int = Query(default=0, ge=0),
    limit: int = Query(default=DEFAULT_CHUNK_SIZE, ge=1),
) -> dict:
    """Read text by byte range with UTF-8 boundary protection."""
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)
    try:
        async with _FILESYSTEM_SEMAPHORE:
            return await asyncio.to_thread(
                read_file_chunk,
                files_root,
                path,
                offset,
                limit,
            )
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except ValueError as exc:
        raise HTTPException(status_code=416, detail=str(exc)) from exc
    except FileVersionConflict as exc:
        raise HTTPException(
            status_code=409,
            detail="File changed while it was being read",
        ) from exc
    except (FileNotFoundError, OSError) as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc


@router.put(
    "/file-content",
    summary="Save workspace text with optimistic concurrency",
)
async def write_workspace_file_content(
    request: Request,
    path: str = Query(...),
    root: str = Query(default="project"),
    body: dict = Body(...),
) -> dict:
    """Atomically save text when the supplied ETag still matches."""
    content = body.get("content")
    if not isinstance(content, str):
        raise HTTPException(status_code=422, detail="content must be a string")
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)
    try:
        async with _FILESYSTEM_SEMAPHORE:
            return await asyncio.to_thread(
                save_text_file,
                files_root,
                path,
                content,
                request.headers.get("if-match"),
            )
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except FileVersionConflict as exc:
        raise HTTPException(
            status_code=409,
            detail="File changed on disk",
        ) from exc
    except OSError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.get(
    "/file-download",
    summary="Stream one workspace file",
)
async def download_workspace_file(
    request: Request,
    path: str = Query(...),
    root: str = Query(default="project"),
) -> StreamingResponse:
    """Stream one safe workspace file without buffering it in memory."""
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)

    def _resolve_download() -> tuple[Path, os.stat_result, str, str]:
        target = resolve_workspace_path(files_root, path)
        info = target.stat()
        filename = target.name.replace('"', "")
        media_type = (
            mimetypes.guess_type(filename)[0] or "application/octet-stream"
        )
        return target, info, filename, media_type

    try:
        async with _FILESYSTEM_SEMAPHORE:
            target, info, filename, media_type = await asyncio.to_thread(
                _resolve_download,
            )
        if not stat.S_ISREG(info.st_mode):
            raise FileNotFoundError(path)
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except (FileNotFoundError, OSError) as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc

    def _stream_file(chunk_size: int = 256 * 1024):
        with target.open("rb") as handle:
            while chunk := handle.read(chunk_size):
                yield chunk

    quoted_filename = quote(filename)
    if quoted_filename == filename:
        content_disposition = f'attachment; filename="{filename}"'
    else:
        content_disposition = f"attachment; filename*=utf-8''{quoted_filename}"
    return StreamingResponse(
        _stream_file(),
        media_type=media_type,
        headers={
            "Accept-Ranges": "bytes",
            "Content-Disposition": content_disposition,
            "Content-Length": str(info.st_size),
            "ETag": file_etag(info),
        },
    )


@router.get(
    "/html-file-uri",
    summary="Resolve one workspace HTML file for the desktop browser",
)
async def resolve_workspace_html_file_uri(
    request: Request,
    path: str = Query(...),
    root: str = Query(default="project"),
) -> dict:
    """Return the URI of one validated HTML file in the selected workspace."""
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)

    def _resolve_html() -> Path:
        target = resolve_workspace_path(files_root, path)
        if target.suffix.lower() not in {".html", ".htm"}:
            raise InvalidWorkspacePath("Path must reference an HTML file")
        if not target.is_file():
            raise FileNotFoundError(path)
        return target

    try:
        async with _FILESYSTEM_SEMAPHORE:
            target = await asyncio.to_thread(_resolve_html)
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except (FileNotFoundError, OSError) as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc

    return {"uri": target.as_uri()}


def _reserve_path(target: Path) -> bool:
    """Atomically reserve one upload target without truncating a file."""
    try:
        descriptor = os.open(
            target,
            os.O_CREAT | os.O_EXCL | os.O_WRONLY,
            0o600,
        )
    except FileExistsError:
        return False
    os.close(descriptor)
    return True


def _reserve_upload_targets(
    upload_targets: list[tuple[UploadFile, str, Path]],
    conflict: str | None,
) -> tuple[list[tuple[UploadFile, str, Path | None, Path]], set[Path]]:
    """Atomically allocate all non-overwrite upload destinations."""
    allocated: list[tuple[UploadFile, str, Path | None, Path]] = []
    reservations: set[Path] = set()
    try:
        for upload, filename, target in upload_targets:
            if conflict == "overwrite":
                allocated.append((upload, filename, target, target))
                continue
            if _reserve_path(target):
                reservations.add(target)
                allocated.append((upload, filename, target, target))
                continue
            if conflict == "skip":
                allocated.append((upload, filename, None, target))
                continue
            if conflict != "rename":
                raise FileExistsError(filename)
            for index in range(1, 10_000):
                candidate = target.with_name(
                    f"{target.stem} ({index}){target.suffix}",
                )
                if _reserve_path(candidate):
                    reservations.add(candidate)
                    allocated.append((upload, filename, candidate, target))
                    break
            else:
                raise OSError("Unable to allocate a conflict-free filename")
    except BaseException:
        for reservation in reservations:
            reservation.unlink(missing_ok=True)
        raise
    return allocated, reservations


def _write_reserved_upload(upload: UploadFile, target: Path) -> int:
    """Copy one upload and atomically replace its reserved target."""
    temporary = target.with_name(
        f".{target.name}.{secrets.token_hex(6)}.qwenpaw.tmp",
    )
    size = 0
    try:
        upload.file.seek(0)
        with temporary.open("wb") as handle:
            while chunk := upload.file.read(256 * 1024):
                size += len(chunk)
                handle.write(chunk)
            handle.flush()
        os.replace(temporary, target)
    except BaseException:
        temporary.unlink(missing_ok=True)
        raise
    return size


def _cleanup_upload_reservations(reservations: set[Path]) -> None:
    """Remove placeholders that were not replaced by completed uploads."""
    for reservation in reservations:
        reservation.unlink(missing_ok=True)


def _filesystem_name_rules(directory: Path) -> tuple[bool, bool]:
    """Detect case and Unicode normalization sensitivity for a directory.

    Thin wrapper over the shared probe: the temp-file technique this used
    to implement inline now lives in
    :mod:`qwenpaw.services.fs_name_rules`, unchanged in behaviour. It stays
    a write probe because the question here is about names that do *not*
    exist yet — would these two uploads collide? — which nothing that
    inspects existing entries can answer.

    Project-directory comparison deliberately does **not** use this. There
    the directories exist, so ``dir_key`` asks which entry each path
    reaches and gets an exact answer; a name-rules guess would be both
    weaker and, for a mount point, wrong.
    """
    rules = probe_name_rules(directory)
    return rules.case_sensitive, rules.normalization_sensitive


def _upload_name_key(
    filename: str,
    *,
    case_sensitive: bool,
    normalization_sensitive: bool,
) -> str:
    """Build a filename comparison key matching the target filesystem."""
    return NameRules(
        case_sensitive=case_sensitive,
        normalization_sensitive=normalization_sensitive,
    ).key(filename)


def _prepare_upload_targets(
    directory: Path,
    files: list[UploadFile],
) -> tuple[list[tuple[UploadFile, str, Path]], list[str]]:
    """Validate upload names and collect conflicts before writing files."""
    upload_targets: list[tuple[UploadFile, str, Path]] = []
    seen_names: set[str] = set()
    conflicts: list[str] = []
    case_sensitive, normalization_sensitive = _filesystem_name_rules(
        directory,
    )
    for upload in files:
        filename = upload.filename or ""
        if "/" in filename or "\\" in filename:
            raise HTTPException(
                status_code=400,
                detail="Upload filename must not contain a path",
            )
        try:
            target = resolve_workspace_path(
                directory,
                filename,
                portable=True,
            )
        except InvalidWorkspacePath as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        comparable_name = _upload_name_key(
            filename,
            case_sensitive=case_sensitive,
            normalization_sensitive=normalization_sensitive,
        )
        if target.exists() or comparable_name in seen_names:
            conflicts.append(filename)
        seen_names.add(comparable_name)
        upload_targets.append((upload, filename, target))
    return upload_targets, conflicts


@router.post(
    "/file-upload",
    summary="Stream ordinary files into one workspace directory",
)
async def upload_workspace_files(
    request: Request,
    files: list[UploadFile] = File(...),
    path: str = Query(default=""),
    root: str = Query(default="project"),
    conflict: str | None = Query(default=None),
) -> dict:
    """Upload files, requesting a policy only when names conflict."""
    if conflict is not None and conflict not in {
        "overwrite",
        "skip",
        "rename",
    }:
        raise HTTPException(
            status_code=400,
            detail="conflict must be overwrite, skip, or rename",
        )
    workspace = await get_agent_for_request(request)
    files_root = await _resolve_files_root(request, workspace, root)

    def _resolve_directory() -> Path:
        directory = resolve_workspace_path(
            files_root,
            path,
            allow_root=True,
        )
        if not directory.is_dir():
            raise NotADirectoryError(path)
        return directory

    try:
        directory = await asyncio.to_thread(_resolve_directory)
    except InvalidWorkspacePath as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except (FileNotFoundError, NotADirectoryError) as exc:
        raise HTTPException(
            status_code=404,
            detail="Upload directory not found",
        ) from exc

    upload_targets, conflicts = await asyncio.to_thread(
        _prepare_upload_targets,
        directory,
        files,
    )

    if conflicts and conflict is None:
        raise HTTPException(
            status_code=409,
            detail={
                "code": "upload_conflict",
                "files": conflicts,
            },
        )

    try:
        allocated, reservations = await asyncio.to_thread(
            _reserve_upload_targets,
            upload_targets,
            conflict,
        )
    except FileExistsError as exc:
        raise HTTPException(
            status_code=409,
            detail={
                "code": "upload_conflict",
                "files": [str(exc)],
            },
        ) from exc

    results: list[dict] = []
    try:
        for upload, filename, target, requested_target in allocated:
            if target is None:
                results.append(
                    {
                        "name": filename,
                        "path": requested_target.relative_to(
                            files_root,
                        ).as_posix(),
                        "status": "skipped",
                    },
                )
                continue
            async with _FILESYSTEM_SEMAPHORE:
                size = await asyncio.to_thread(
                    _write_reserved_upload,
                    upload,
                    target,
                )
                reservations.discard(target)

            results.append(
                {
                    "name": filename,
                    "path": target.relative_to(files_root).as_posix(),
                    "size": size,
                    "status": "uploaded",
                },
            )
    finally:
        await asyncio.to_thread(
            _cleanup_upload_reservations,
            reservations,
        )
    return {"files": results}


@router.get(
    "/code-files",
    summary="List all workspace files (Coding Mode)",
)
async def list_code_files(request: Request) -> list[dict]:
    """List every non-hidden file in the active coding project directory."""
    workspace = await get_agent_for_request(request)
    return await asyncio.to_thread(
        lambda: _list_all_files(get_agent_project_dir(workspace)),
    )


_CODE_FILE_MAX_BYTES = 5 * 1024 * 1024  # 5 MB
_BINARY_FILE_MAX_BYTES = 50 * 1024 * 1024  # 50 MB
_IMAGE_FILE_MAX_BYTES = 100 * 1024 * 1024  # 100 MB
_DOCUMENT_FILE_MAX_BYTES = 200 * 1024 * 1024  # 200 MB
_MEDIA_FILE_MAX_BYTES = 2 * 1024 * 1024 * 1024  # 2 GB

_MIME_MAP: dict[str, str] = {
    # Images
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
    "svg": "image/svg+xml",
    "ico": "image/x-icon",
    "bmp": "image/bmp",
    "tiff": "image/tiff",
    "tif": "image/tiff",
    # Documents
    "pdf": "application/pdf",
    # Office (served as binary for download/preview)
    "doc": "application/msword",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # noqa: E501
    "xls": "application/vnd.ms-excel",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # noqa: E501
    "ppt": "application/vnd.ms-powerpoint",
    "pptx": (
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation"
    ),
    "odt": "application/vnd.oasis.opendocument.text",
    "ods": "application/vnd.oasis.opendocument.spreadsheet",
    "odp": "application/vnd.oasis.opendocument.presentation",
    # Video
    "mp4": "video/mp4",
    "webm": "video/webm",
    "avi": "video/x-msvideo",
    "mov": "video/quicktime",
    "mkv": "video/x-matroska",
    "wmv": "video/x-ms-wmv",
    "flv": "video/x-flv",
    # Audio
    "mp3": "audio/mpeg",
    "wav": "audio/wav",
    "flac": "audio/flac",
    "aac": "audio/aac",
    "ogg": "audio/ogg",
    "wma": "audio/x-ms-wma",
    # Archives
    "zip": "application/zip",
    "tar": "application/x-tar",
    "gz": "application/gzip",
    "7z": "application/x-7z-compressed",
    "rar": "application/vnd.rar",
    # Data
    "csv": "text/csv",
}


def _binary_file_size_limit(mime: str) -> int:
    """Return the preview limit for the resource category."""
    if mime.startswith(("video/", "audio/")):
        return _MEDIA_FILE_MAX_BYTES
    if mime.startswith("image/"):
        return _IMAGE_FILE_MAX_BYTES
    if mime == "application/pdf" or mime.startswith("application/vnd."):
        return _DOCUMENT_FILE_MAX_BYTES
    return _BINARY_FILE_MAX_BYTES


def _parse_single_byte_range(value: str, size: int) -> tuple[int, int]:
    """Parse one RFC 7233 byte range and return inclusive bounds."""
    if size <= 0 or not value.startswith("bytes="):
        raise ValueError("Invalid byte range")

    spec = value[6:].strip()
    if not spec or "," in spec or "-" not in spec:
        raise ValueError("Only a single byte range is supported")

    start_text, end_text = (part.strip() for part in spec.split("-", 1))
    if not start_text:
        if not end_text.isdigit():
            raise ValueError("Invalid suffix byte range")
        suffix_length = int(end_text)
        if suffix_length <= 0:
            raise ValueError("Invalid suffix byte range")
        start = max(size - suffix_length, 0)
        return start, size - 1

    if not start_text.isdigit() or (end_text and not end_text.isdigit()):
        raise ValueError("Invalid byte range")

    start = int(start_text)
    if start >= size:
        raise ValueError("Byte range starts beyond end of file")

    end = int(end_text) if end_text else size - 1
    if end < start:
        raise ValueError("Byte range end precedes start")
    return start, min(end, size - 1)


@router.get(
    "/binary-files/{file_path:path}",
    summary="Serve a binary workspace file (images, PDFs, CSV) for preview",
)
async def read_binary_file(  # pylint: disable=too-many-statements
    # pylint: disable=unused-argument
    file_path: str,
    request: Request,
    agent_id: str | None = None,
) -> StreamingResponse:
    """Return the raw bytes of *file_path* with the appropriate Content-Type.

    Intended for the IDE preview panel (images, PDFs, CSV).
    Accepts relative paths within the workspace or absolute paths
    (from ``file://`` URLs produced by tool-call results).
    Supports one RFC 7233 byte range for seekable media playback. Size limits
    vary by resource category; requests without Range still stream the full
    file for clients that do not implement partial loading.
    """
    workspace = await get_agent_for_request(request)
    target = await asyncio.to_thread(
        lambda: safe_join(get_agent_project_dir(workspace), file_path),
    )

    ext = target.suffix.lstrip(".").lower()
    mime = _MIME_MAP.get(ext)
    if mime is None:
        raise HTTPException(
            status_code=415,
            detail=f"Preview not supported for .{ext} files",
        )

    try:
        size = await asyncio.to_thread(lambda: target.stat().st_size)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc
    except OSError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    size_limit = _binary_file_size_limit(mime)
    if size > size_limit:
        raise HTTPException(
            status_code=413,
            detail=(
                f"File too large for preview ({size // 1024 // 1024} MB"
                f" > {size_limit // 1024 // 1024} MB limit)"
            ),
        )

    range_header = request.headers.get("range")
    start = 0
    end = size - 1
    status_code = 200
    headers = {
        "Accept-Ranges": "bytes",
        "Content-Length": str(size),
        # Explicitly tell the browser/WebView to display inline rather
        # than triggering a download dialog (especially for PDFs).
        "Content-Disposition": "inline",
    }
    if range_header is not None:
        try:
            start, end = _parse_single_byte_range(range_header, size)
        except ValueError as exc:
            raise HTTPException(
                status_code=416,
                detail=str(exc),
                headers={
                    "Accept-Ranges": "bytes",
                    "Content-Range": f"bytes */{size}",
                },
            ) from exc
        status_code = 206
        headers["Content-Length"] = str(end - start + 1)
        headers["Content-Range"] = f"bytes {start}-{end}/{size}"

    chunk_size = (
        1024 * 1024 if mime.startswith(("video/", "audio/")) else 64 * 1024
    )

    def _iter_chunks():
        remaining = end - start + 1
        with open(target, "rb") as fh:
            fh.seek(start)
            while remaining > 0:
                data = fh.read(min(chunk_size, remaining))
                if not data:
                    break
                remaining -= len(data)
                yield data

    return StreamingResponse(
        _iter_chunks(),
        status_code=status_code,
        media_type=mime,
        headers=headers,
    )


def _file_etag(stat_result: os.stat_result) -> str:
    """Build a weak ETag from mtime+size — cheap and good enough for IDE."""
    return f'W/"{stat_result.st_mtime_ns}-{stat_result.st_size}"'


@router.get(
    "/code-files/{file_path:path}",
    summary="Read any workspace file (Coding Mode)",
)
async def read_code_file(file_path: str, request: Request):
    """Return the text content of *file_path* inside the workspace.

    Adds a weak ETag (mtime_ns + size) so repeat opens of an unchanged file
    short-circuit to ``304 Not Modified`` and skip the read entirely.
    Returns HTTP 413 if the file exceeds ``_CODE_FILE_MAX_BYTES`` (5 MB) to
    avoid flooding the browser with huge binary or log files.
    """
    workspace = await get_agent_for_request(request)
    target = await asyncio.to_thread(
        lambda: safe_join(get_agent_project_dir(workspace), file_path),
    )

    def _stat() -> os.stat_result:
        return target.stat()

    try:
        st = await asyncio.to_thread(_stat)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail="File not found") from exc
    except OSError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    if not stat.S_ISREG(st.st_mode):
        raise HTTPException(status_code=404, detail="File not found")

    etag = _file_etag(st)
    if request.headers.get("if-none-match") == etag:
        return Response(status_code=304, headers={"ETag": etag})

    if st.st_size > _CODE_FILE_MAX_BYTES:
        raise HTTPException(
            status_code=413,
            detail=(
                f"File too large to open in editor "
                f"({st.st_size // 1024 // 1024} MB"
                f" > {_CODE_FILE_MAX_BYTES // 1024 // 1024} MB limit)"
            ),
        )

    def _read() -> str:
        return target.read_text(encoding="utf-8", errors="replace")

    try:
        content = await asyncio.to_thread(_read)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    return ORJSONResponse(
        {"path": file_path, "content": content},
        headers={"ETag": etag},
    )


@router.put(
    "/code-files/{file_path:path}",
    summary="Write any workspace file (Coding Mode)",
)
async def write_code_file(
    file_path: str,
    request: Request,
    body: dict = Body(...),
) -> dict:
    """Overwrite *file_path* inside the workspace with the provided content.

    Request body::

        {"content": "<new file content>"}
    """
    workspace = await get_agent_for_request(request)
    target = await asyncio.to_thread(
        lambda: safe_join(get_agent_project_dir(workspace), file_path),
    )
    content = body.get("content", "")
    if not isinstance(content, str):
        raise HTTPException(status_code=422, detail="content must be a string")

    def _write() -> int:
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(content.encode("utf-8"))
        return target.stat().st_size

    try:
        size = await asyncio.to_thread(_write)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    return {"path": file_path, "size": size}


# ---------------------------------------------------------------------------
# Office document conversion (DOCX → HTML)
# ---------------------------------------------------------------------------


class ConvertOfficeRequest(BaseModel):
    """Request body for /convert-office."""

    url: str = Field(..., description="File URL or path")
    mime_type: str | None = Field(None, description="MIME type of the file")


class OfficeViewRequest(BaseModel):
    """Request body for /office-screenshot and /office-outline."""

    url: str = Field(..., description="File URL or path")
    page: int = Field(1, description="Page/slide number (1-based)")


# ---------------------------------------------------------------------------
# OfficeCLI integration helpers
# ---------------------------------------------------------------------------

_OFFICECLI_TIMEOUT = 30  # seconds

# Cache for officecli availability check (avoids repeated subprocess calls)
_officecli_checked: bool = False
_officecli_ok: bool = False


def _bundled_officecli_path() -> str | None:
    """Return the path to a bundled officecli binary, if available.

    In the Tauri desktop build, the officecli binary is shipped as a
    resource under ``binaries/officecli/``. The Rust backend launcher
    sets ``QWENPAW_DESKTOP_OFFICECLI_DIR`` to that directory.

    Returns the full path to the executable, or ``None`` if not found.
    """
    oc_dir = os.environ.get("QWENPAW_DESKTOP_OFFICECLI_DIR")
    if not oc_dir:
        return None
    exe_name = "officecli.exe" if sys.platform == "win32" else "officecli"
    candidate = Path(oc_dir) / exe_name
    if candidate.is_file():
        return str(candidate)
    return None


def _officecli_bin() -> str:
    """Return the resolved officecli executable path.

    Priority:
    1. Bundled binary in the Tauri resource directory (desktop app)
    2. ``shutil.which("officecli")`` (system PATH — npm install, etc.)
    3. Bare ``"officecli"`` as a last resort (will likely fail)
    """
    bundled = _bundled_officecli_path()
    if bundled:
        return bundled
    resolved = shutil.which("officecli")
    return resolved or "officecli"


def _officecli_env() -> dict[str, str]:
    """Build environment variables for officecli subprocess calls.

    When using a bundled officecli, the licensed document-reader plugin
    (``plugins/dump-reader/doc/plugin``) must be exported via an
    environment variable so officecli can find it at runtime.
    """
    env: dict[str, str] = {}
    bundled = _bundled_officecli_path()
    if not bundled:
        return env
    oc_dir = Path(bundled).parent
    plugin = oc_dir / "plugins" / "dump-reader" / "doc" / "plugin"
    if plugin.is_file():
        env["OFFICECLI_PLUGIN_DUMP_READER_DOC"] = str(plugin)
    return env


def _is_officecli_available() -> bool:
    """Check if the correct officecli binary is available and supports
    ``view``.

    The npm package ``officecli`` (v0.2.x) is a *different* tool — an AI
    document generator that does NOT have a ``view`` subcommand.  This
    function verifies that the installed binary is the OfficeCLI from
    https://github.com/iOfficeAI/OfficeCLI/releases which supports
    ``view <file> html``, ``view <file> screenshot``, etc.

    The check is cached after the first call to avoid repeated subprocess
    invocations.
    """
    global _officecli_checked, _officecli_ok  # noqa: PLW0603
    if _officecli_checked:
        return _officecli_ok

    _officecli_checked = True
    resolved = _bundled_officecli_path() or shutil.which("officecli")
    if resolved is None:
        _officecli_ok = False
        return False

    # Verify it supports the "view" subcommand by checking --help output.
    # The wrong officecli (npm AI generator) will either error out or
    # not list "view" in its commands.
    try:
        result = subprocess.run(  # noqa: S603
            [resolved, "--help"],
            capture_output=True,
            text=True,
            timeout=5,
            check=False,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        _officecli_ok = False
        return False

    if result.returncode != 0:
        # Some CLIs print help to stderr with non-zero exit
        help_text = (result.stdout or "") + (result.stderr or "")
    else:
        help_text = result.stdout or ""

    # The correct OfficeCLI has "view" in its help/commands list.
    # The npm officecli has "new", "doctor", "login", etc. but NOT "view".
    _officecli_ok = "view" in help_text
    if not _officecli_ok:
        logger.warning(
            "officecli binary found at %s but does not support 'view' "
            "subcommand — likely the wrong npm package. "
            "Install the correct one from "
            "https://github.com/iOfficeAI/OfficeCLI/releases",
            resolved,
        )
    else:
        logger.info("officecli (with view support) detected at %s", resolved)
    return _officecli_ok


def _convert_with_officecli(file_path: str) -> str | None:
    """High-fidelity HTML conversion via ``officecli view <file> html``.

    Reads HTML directly from stdout (no temp file needed), which is
    significantly faster than the ``-o <tmpfile>`` approach because it
    avoids filesystem I/O.
    Returns an HTML string on success, or *None* on any failure so the
    caller can fall back to the legacy conversion path.
    """
    try:
        result = subprocess.run(  # noqa: S603
            [
                _officecli_bin(),
                "view",
                file_path,
                "html",
            ],
            capture_output=True,
            timeout=_OFFICECLI_TIMEOUT,
            check=False,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
        logger.warning(
            "officecli conversion failed for %s: %s",
            file_path,
            exc,
        )
        return None
    if result.returncode != 0:
        stderr = (
            result.stderr.decode(errors="replace") if result.stderr else ""
        )
        logger.warning(
            "officecli returned non-zero exit %d for %s: %s",
            result.returncode,
            file_path,
            stderr[:200],
        )
        return None
    html = result.stdout.decode(errors="replace") if result.stdout else ""
    if not html.strip():
        logger.warning("officecli returned empty HTML for %s", file_path)
        return None
    return html


_page_count_cache: dict[str, tuple[int, float]] = {}
_PAGE_COUNT_CACHE_TTL = 300.0  # 5 minutes
_PAGE_COUNT_CACHE_MAX_SIZE = 100


def _get_officecli_page_count(file_path: str) -> int:
    """Get the total page/slide count of an Office document via officecli.

    Uses ``officecli view <file> stats --json``.
    Results are cached for 5 minutes per file to avoid repeated subprocess
    invocations on every screenshot request.
    Returns 0 if the count cannot be determined (best-effort).
    """
    import time as _time

    now = _time.time()

    # Evict expired entries to prevent unbounded growth.
    expired = [
        k
        for k, (_, ts) in _page_count_cache.items()
        if now - ts >= _PAGE_COUNT_CACHE_TTL
    ]
    for k in expired:
        del _page_count_cache[k]

    # If still over the size cap, drop the oldest entries.
    if len(_page_count_cache) > _PAGE_COUNT_CACHE_MAX_SIZE:
        for k, _ in sorted(
            _page_count_cache.items(),
            key=lambda item: item[1][1],
        )[: len(_page_count_cache) - _PAGE_COUNT_CACHE_MAX_SIZE]:
            del _page_count_cache[k]

    cached = _page_count_cache.get(file_path)
    if cached and (now - cached[1]) < _PAGE_COUNT_CACHE_TTL:
        return cached[0]

    try:
        result = subprocess.run(  # noqa: S603
            [
                _officecli_bin(),
                "view",
                file_path,
                "stats",
                "--json",
            ],
            capture_output=True,
            text=True,
            timeout=_OFFICECLI_TIMEOUT,
            check=False,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return 0
    if result.returncode != 0:
        return 0
    try:
        data = json.loads(result.stdout)
    except json.JSONDecodeError:
        return 0
    # officecli stats --json returns:
    #   {"success": true, "data": {"slides": 5, ...}}   (PPTX)
    #   {"success": true, "data": {"pageCount": 3, ...}} (DOCX)
    # Try top-level and nested "data" keys
    nested = data.get("data") or {}
    count = (
        data.get("pageCount")
        or data.get("page_count")
        or data.get("slides")
        or data.get("totalPages")
        or nested.get("pageCount")
        or nested.get("page_count")
        or nested.get("slides")
        or nested.get("totalPages")
        or 0
    )
    _page_count_cache[file_path] = (count, now)
    return count


def _ensure_path_in_allowed_roots(
    target: Path,
    coding_dir: Path,
    workspace_dir: Path | None = None,
) -> Path:
    """Ensure *target* resolves within coding_dir, workspace_dir,
    or WORKING_DIR.

    Absolute paths (e.g. from ``file://`` URLs produced by tool-call
    results) are accepted only when they resolve within the coding
    project directory, the agent workspace, or the global QwenPaw
    working directory.  This mirrors the Tauri backend's
    ``resolve_workspace_file_path`` dual-root check and prevents
    arbitrary filesystem reads via crafted absolute paths.

    Raises ``HTTPException(403)`` when the path is outside all roots.
    """
    resolved = target.resolve()
    allowed_roots: list[Path] = []
    for root in (coding_dir, workspace_dir, WORKING_DIR):
        if root is not None:
            try:
                allowed_roots.append(root.resolve())
            except OSError:
                pass
    for root in allowed_roots:
        try:
            resolved.relative_to(root)
            return resolved
        except ValueError:
            continue
    raise HTTPException(
        status_code=403,
        detail="Path is outside the allowed workspace directories",
    )


def _resolve_file_path_from_url(  # pylint: disable=too-many-branches
    # pylint: disable=too-many-statements
    url: str,
    coding_dir: Path,
    workspace_dir: Path | None = None,
) -> Path:
    """Resolve a frontend URL to an absolute file path.

    Handles ``/api/workspace/binary-files/<path>``, absolute paths, and
    ``file://`` URLs.  Raises ``HTTPException(404)`` if the file does
    not exist, or ``HTTPException(403)`` if the resolved path is outside
    the allowed workspace roots.
    """
    parsed_url = urlparse(url)
    if parsed_url.path.endswith("/workspace/file-download"):
        query = parse_qs(parsed_url.query)
        paths = query.get("path")
        if not paths or not paths[0]:
            raise HTTPException(
                status_code=422,
                detail="Workspace file-download URL is missing path",
            )
        root = query.get("root", ["project"])[0]
        if root == "workspace":
            if workspace_dir is None:
                raise HTTPException(
                    status_code=404,
                    detail="Agent workspace directory is unavailable",
                )
            base_dir = workspace_dir
        elif root == "project":
            base_dir = coding_dir
        else:
            raise HTTPException(
                status_code=400,
                detail="root must be project or workspace",
            )
        try:
            target = safe_join(base_dir, paths[0])
        except Exception as exc:
            raise HTTPException(
                status_code=403,
                detail="Path is outside the allowed workspace directories",
            ) from exc
        if not target.exists():
            raise HTTPException(
                status_code=404,
                detail=f"File not found: {paths[0]}",
            )
        return target
    elif "/files/preview/" in parsed_url.path:
        encoded_path = parsed_url.path.split("/files/preview/", 1)[1]
        file_path = unquote(encoded_path)
    elif "/binary-files/" in url:
        file_path = url.split("/binary-files/", 1)[1].split("?")[0]
        file_path = unquote(file_path)
    elif url.startswith("/"):
        file_path = url.lstrip("/")
    elif url.startswith("file://"):
        # file:///tmp/test.docx → /tmp/test.docx (Unix absolute)
        # file:///C:/Users/...  → C:/Users/...  (Windows drive letter)
        # file://localhost/path → /path          (localhost authority)
        # file://host/share/... → //host/share/  (UNC)
        rest = url[len("file://") :]
        if rest.startswith("/"):
            # file:///path → /path (Unix) or file:///C:/... → /C:/... (Windows)
            file_path = unquote(rest)
            # On Windows, strip leading "/" before a drive letter: /C:/x → C:/x
            if (
                len(file_path) > 2
                and file_path[0] == "/"
                and file_path[2] == ":"
                and file_path[1].isalpha()
            ):
                file_path = file_path[1:]
        elif rest.startswith("localhost/"):
            # file://localhost/path → /path
            file_path = unquote(rest[len("localhost") :])
        else:
            # file://host/share/... → //host/share/... (UNC)
            file_path = unquote("//" + rest)
    else:
        file_path = url

    # If the extracted path is absolute (common with file:// URLs from
    # tool-call results), verify it is within the allowed workspace roots
    # before using it.  Relative paths go through safe_join for the same
    # containment guarantee.
    extracted = Path(file_path)
    if extracted.is_absolute():
        target = _ensure_path_in_allowed_roots(
            extracted,
            coding_dir,
            workspace_dir,
        )
    else:
        try:
            target = safe_join(coding_dir, file_path)
        except Exception:
            target = _ensure_path_in_allowed_roots(
                extracted,
                coding_dir,
                workspace_dir,
            )

    if not target.exists():
        raise HTTPException(
            status_code=404,
            detail=f"File not found: {file_path}",
        )
    return target


def _get_docx_page_info(file_path: str) -> dict:
    """Parse DOCX XML to get page dimensions and estimate lines/chars."""
    import xml.etree.ElementTree as ET

    W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    # Defaults: US Letter, 1-inch margins
    defaults = {
        "lines_per_page": 45,
        "chars_per_line": 78,
    }

    try:
        with zipfile.ZipFile(file_path) as z:
            root = ET.fromstring(z.read("word/document.xml"))
    except Exception:
        return defaults

    body = root.find(f"{{{W_NS}}}body")
    if body is None:
        return defaults

    # Find section properties (document-level)
    sect_pr = body.find(f"{{{W_NS}}}sectPr")
    if sect_pr is None:
        for child in body:
            pPr = child.find(f"{{{W_NS}}}pPr")
            if pPr is not None:
                s = pPr.find(f"{{{W_NS}}}sectPr")
                if s is not None:
                    sect_pr = s
                    break
    if sect_pr is None:
        return defaults

    pg_sz = sect_pr.find(f"{{{W_NS}}}pgSz")
    pg_mar = sect_pr.find(f"{{{W_NS}}}pgMar")

    page_w = 12240
    page_h = 15840
    margin_top = 1440
    margin_bottom = 1440
    margin_left = 1440
    margin_right = 1440
    header_h = 0
    footer_h = 0

    if pg_sz is not None:
        page_w = int(pg_sz.get(f"{{{W_NS}}}w", page_w))
        page_h = int(pg_sz.get(f"{{{W_NS}}}h", page_h))
    if pg_mar is not None:
        margin_top = int(pg_mar.get(f"{{{W_NS}}}top", margin_top))
        margin_bottom = int(pg_mar.get(f"{{{W_NS}}}bottom", margin_bottom))
        margin_left = int(pg_mar.get(f"{{{W_NS}}}left", margin_left))
        margin_right = int(pg_mar.get(f"{{{W_NS}}}right", margin_right))
        header_h = int(pg_mar.get(f"{{{W_NS}}}header", 0))
        footer_h = int(pg_mar.get(f"{{{W_NS}}}footer", 0))

    usable_h = (
        page_h
        - margin_top
        - margin_bottom
        - max(0, header_h)
        - max(0, footer_h)
    )
    usable_w = page_w - margin_left - margin_right

    # 11pt font, 1.15 line spacing → ~253 twips/line
    lines_per_page = max(10, usable_h // 253)
    # 11pt font → avg char width ~120 twips
    chars_per_line = max(20, usable_w // 120)

    return {
        "lines_per_page": lines_per_page,
        "chars_per_line": chars_per_line,
    }


def _has_page_break(element) -> bool:
    """Recursively check if an element tree contains a page break."""
    try:
        from mammoth import documents
    except ImportError:
        return False
    if isinstance(element, documents.Break) and element.break_type == "page":
        return True
    if hasattr(element, "children"):
        return any(_has_page_break(c) for c in element.children)
    return False


def _estimate_element_lines(  # pylint: disable=too-many-return-statements
    element,
    chars_per_line: int,
) -> float:
    """Estimate the number of visual lines an element occupies."""
    try:
        from mammoth import documents
    except ImportError:
        return 1.0

    if isinstance(element, documents.Paragraph):
        # Collect text content
        text = []

        def collect_text(el):
            if isinstance(el, documents.Text):
                text.append(el.value)
            elif hasattr(el, "children"):
                for c in el.children:
                    collect_text(c)

        collect_text(element)
        full_text = "".join(text)

        if not full_text.strip():
            return 1.0

        text_lines = max(1.0, len(full_text) / chars_per_line)

        style_name = (element.style_name or "").lower()
        if "heading 1" in style_name:
            return text_lines + 2.0
        elif "heading 2" in style_name:
            return text_lines + 1.5
        elif "heading 3" in style_name:
            return text_lines + 1.0
        elif "title" in style_name:
            return text_lines + 3.0
        return text_lines

    elif isinstance(element, documents.Table):
        row_count = len(element.children)
        return max(3.0, row_count * 1.5)

    elif hasattr(element, "children"):
        total = 0.0
        for child in element.children:
            total += _estimate_element_lines(child, chars_per_line)
        return max(1.0, total)

    return 1.0


# Sentinel text used as a page-break marker inside mammoth's HTML output.
_PAGE_BREAK_MARKER = "\u0000PAGE_BREAK\u0000"


def _build_docx_transform(page_info: dict):
    """Build a transform_document function that inserts page-break markers."""
    try:
        from mammoth import documents
    except ImportError:
        return None

    lines_per_page = page_info.get("lines_per_page", 45)
    chars_per_line = page_info.get("chars_per_line", 78)

    def transform_document(document):
        new_children = []
        accumulated_lines = 0.0

        for child in document.children:
            has_explicit = _has_page_break(child)
            element_lines = _estimate_element_lines(child, chars_per_line)

            needs_break = False
            if has_explicit:
                needs_break = True
                accumulated_lines = element_lines
            elif (
                accumulated_lines > 0
                and accumulated_lines + element_lines > lines_per_page
            ):
                needs_break = True
                accumulated_lines = element_lines
            else:
                accumulated_lines += element_lines

            if needs_break:
                marker_run = documents.run(
                    children=[documents.text(_PAGE_BREAK_MARKER)],
                )
                marker_para = documents.paragraph(
                    style_id="PageBreakMarker",
                    style_name="Page Break Marker",
                    numbering=None,
                    alignment=None,
                    indent=None,
                    children=[marker_run],
                )
                new_children.append(marker_para)

            new_children.append(child)

        return document.copy(children=new_children)

    return transform_document


def _html_escape(text: str) -> str:
    """Escape HTML special characters to prevent injection in legacy
    conversion."""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&#x27;")
    )


def _convert_legacy_office_to_html(file_path: str) -> str:
    """Convert a legacy Office file (.doc, .xls, .ppt) to HTML.

    Uses LibreOffice (``soffice``) to convert to the modern OOXML
    format first, then delegates to :func:`_convert_docx_to_html`
    for HTML rendering.
    """
    source = Path(file_path)
    ext = source.suffix.lower()

    # Map legacy extensions to modern equivalents
    modern_ext = {
        ".doc": ".docx",
        ".xls": ".xlsx",
        ".ppt": ".pptx",
    }.get(ext)
    if modern_ext is None:
        # Not a legacy format — try direct conversion
        return _convert_docx_to_html(file_path)

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        logger.warning(
            "LibreOffice (soffice) not found; cannot convert legacy "
            "format %s",
            file_path,
        )
        return _convert_docx_to_html(file_path)

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            subprocess.run(  # noqa: S603
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    modern_ext.lstrip("."),
                    "--outdir",
                    tmpdir,
                    file_path,
                ],
                capture_output=True,
                timeout=60,
                check=False,
            )
        except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
            logger.warning(
                "LibreOffice conversion failed for %s: %s",
                file_path,
                exc,
            )
            return _convert_docx_to_html(file_path)

        converted = Path(tmpdir) / (source.stem + modern_ext)
        if not converted.is_file():
            logger.warning(
                "LibreOffice did not produce expected output %s",
                converted,
            )
            return _convert_docx_to_html(file_path)

        return _convert_docx_to_html(str(converted))


def _convert_docx_to_html(  # pylint: disable=R0912,R0915
    file_path: str,
) -> str:
    """Convert a .docx/.xlsx/.pptx file to HTML for preview (legacy fallback).

    Uses python libraries (mammoth, openpyxl, python-pptx) for conversion.
    The caller (:func:`convert_office`) is responsible for trying officecli
    first; this function is only the low-fidelity fallback.
    """
    ext = Path(file_path).suffix.lstrip(".").lower()

    if ext == "docx":
        try:
            import mammoth

            # Parse page dimensions for page-break estimation
            page_info = _get_docx_page_info(file_path)
            transform = _build_docx_transform(page_info)

            with open(file_path, "rb") as f:
                if transform:
                    result = mammoth.convert_to_html(
                        f,
                        transform_document=transform,
                    )
                else:
                    result = mammoth.convert_to_html(f)

            html: str = result.value or ""

            # Replace markers with styled page-break divs
            if _PAGE_BREAK_MARKER in html:
                html = html.replace(
                    f"<p>{_PAGE_BREAK_MARKER}</p>",
                    '<div class="docx-page-break"></div>',
                )
                # Clean up any remaining marker text
                html = html.replace(_PAGE_BREAK_MARKER, "")

            return html

        except ImportError:
            pass

        # Fallback: python-docx (basic text extraction)
        try:
            from docx import Document

            doc = Document(file_path)
            html_parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    html_parts.append("<br/>")
                    continue
                style = (para.style.name or "").lower()
                if "heading 1" in style:
                    html_parts.append(f"<h1>{_html_escape(text)}</h1>")
                elif "heading 2" in style:
                    html_parts.append(f"<h2>{_html_escape(text)}</h2>")
                elif "heading 3" in style:
                    html_parts.append(f"<h3>{_html_escape(text)}</h3>")
                elif "title" in style:
                    html_parts.append(
                        f"<h1 style='text-align:center'>"
                        f"{_html_escape(text)}</h1>",
                    )
                else:
                    html_parts.append(f"<p>{_html_escape(text)}</p>")
            # Tables
            for table in doc.tables:
                html_parts.append(
                    "<table border='1' style='border-collapse:collapse'>",
                )
                for row in table.rows:
                    html_parts.append("<tr>")
                    for cell in row.cells:
                        html_parts.append(
                            f"<td>{_html_escape(cell.text)}</td>",
                        )
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            return "\n".join(html_parts)
        except ImportError as exc:
            raise HTTPException(
                status_code=500,
                detail="No docx conversion library available "  # noqa: E501
                "(mammoth or python-docx required)",
            ) from exc

    # For .doc, .xls, .ppt — try LibreOffice if available
    if ext in ("doc", "xls", "ppt", "odt", "ods", "odp"):
        raise HTTPException(
            status_code=415,
            detail=(
                f"Direct preview of .{ext} files is not supported. "
                f"Please convert to .docx/.xlsx/.pptx first."
            ),
        )

    # For .xlsx — basic table extraction
    if ext == "xlsx":
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, read_only=True)
            html_parts = []
            for ws in wb.worksheets:
                html_parts.append(
                    f"<h3>Sheet: {_html_escape(ws.title)}</h3>",
                )
                html_parts.append(
                    "<table border='1' style='border-collapse:collapse'>",
                )
                for row in ws.iter_rows(max_row=100, values_only=True):
                    html_parts.append("<tr>")
                    for cell in row:
                        cell_str = str(cell) if cell is not None else ""
                        html_parts.append(
                            f"<td>{_html_escape(cell_str)}</td>",
                        )
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            return "\n".join(html_parts)
        except ImportError as exc:
            raise HTTPException(
                status_code=500,
                detail="openpyxl not installed for .xlsx preview",
            ) from exc
        except Exception as exc:
            logger.exception("XLSX conversion failed for %s", file_path)
            raise HTTPException(
                status_code=500,
                detail=f"XLSX preview error: {exc}",
            ) from exc

    # For .pptx — basic slide extraction
    if ext == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            html_parts = []
            for i, slide in enumerate(prs.slides, 1):
                html_parts.append(f"<h3>Slide {i}</h3>")
                html_parts.append("<div>")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        html_parts.append(
                            f"<p>{_html_escape(shape.text)}</p>",
                        )
                html_parts.append("</div>")
            return "\n".join(html_parts)
        except ImportError as exc:
            raise HTTPException(
                status_code=500,
                detail="python-pptx not installed for .pptx preview",
            ) from exc
        except Exception as exc:
            logger.exception("PPTX conversion failed for %s", file_path)
            raise HTTPException(
                status_code=500,
                detail=f"PPTX preview error: {exc}",
            ) from exc

    raise HTTPException(
        status_code=415,
        detail=f"Unsupported file type: .{ext}",
    )


@router.post(
    "/convert-office",
    summary="Convert an Office document to HTML for preview",
)
async def convert_office(
    request: Request,
    body: ConvertOfficeRequest,
) -> dict:
    """Convert a .docx/.xlsx/.pptx file to HTML for in-browser preview.

    Accepts a file URL (from the binary-files endpoint) or a workspace
    file path. Uses officecli (high fidelity) when available, falling
    back to mammoth for DOCX, openpyxl for XLSX, python-pptx for PPTX.
    """
    workspace = await get_agent_for_request(request)
    coding_dir = await get_project_dir_for_request(request, workspace)

    target = _resolve_file_path_from_url(
        body.url or "",
        coding_dir,
        workspace.workspace_dir,
    )

    def _convert() -> tuple[str, str]:
        # Try officecli first for high-fidelity rendering
        if _is_officecli_available():
            html = _convert_with_officecli(str(target))
            if html:
                logger.info("officecli conversion succeeded for %s", target)
                return html, "officecli"
            logger.info(
                "officecli conversion returned None, falling back"
                " to legacy for %s",
                target,
            )
        else:
            logger.debug(
                "officecli not available, using legacy conversion for %s",
                target,
            )
        # Fallback to legacy conversion
        return _convert_docx_to_html(str(target)), "legacy"

    try:
        html, engine = await asyncio.to_thread(_convert)
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("convert-office failed for %s", target)
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    return {"html": html, "engine": engine}


@router.post(
    "/office-screenshot",
    summary="Render an Office document page as PNG screenshot",
)
async def office_screenshot(
    request: Request,
    body: OfficeViewRequest,
) -> Response:
    """Render a specific page/slide of an Office document as a PNG image.

    Uses ``officecli view <file> screenshot --page N -o <tmp>``.
    Returns the PNG file as a ``FileResponse``.

    Raises ``HTTPException(404)`` if officecli is not installed.
    """
    if not _is_officecli_available():
        raise HTTPException(
            status_code=404,
            detail="officecli is not installed. "
            "Install from: https://github.com/iOfficeAI/OfficeCLI/releases",
        )

    workspace = await get_agent_for_request(request)
    coding_dir = await get_project_dir_for_request(request, workspace)
    target = _resolve_file_path_from_url(
        body.url,
        coding_dir,
        workspace.workspace_dir,
    )

    # Use mkstemp instead of deprecated mktemp to avoid TOCTOU race.
    # officecli will overwrite the empty file created here.
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".png")
    os.close(tmp_fd)  # Release the fd; officecli opens it via -o

    def _run_screenshot() -> subprocess.CompletedProcess:
        return subprocess.run(  # noqa: S603
            [
                _officecli_bin(),
                "view",
                str(target),
                "screenshot",
                "--page",
                str(body.page),
                "-o",
                tmp_path,
            ],
            capture_output=True,
            timeout=60,
            check=False,
        )

    try:
        result = await asyncio.to_thread(_run_screenshot)
    except subprocess.TimeoutExpired as exc:
        Path(tmp_path).unlink(missing_ok=True)
        raise HTTPException(
            status_code=504,
            detail="officecli screenshot timed out",
        ) from exc

    if result.returncode != 0 or not Path(tmp_path).exists():
        Path(tmp_path).unlink(missing_ok=True)
        stderr = (
            result.stderr.decode(
                errors="replace",
            )
            if result.stderr
            else ""
        )
        raise HTTPException(
            status_code=500,
            detail=f"officecli screenshot failed: {stderr}",
        )

    # Get page count for frontend pagination (best-effort)
    page_count = _get_officecli_page_count(str(target))

    from starlette.background import BackgroundTask

    return FileResponse(
        tmp_path,
        media_type="image/png",
        filename=f"page_{body.page}.png",
        headers={"X-Total-Pages": str(page_count)} if page_count > 0 else {},
        background=BackgroundTask(
            lambda: Path(tmp_path).unlink(missing_ok=True),
        ),
    )


@router.post(
    "/office-outline",
    summary="Get document outline (headings/structure)",
)
async def office_outline(
    request: Request,
    body: ConvertOfficeRequest,
) -> dict:
    """Get the outline (heading structure) of an Office document.

    Uses ``officecli view <file> outline --json``.
    """
    if not _is_officecli_available():
        raise HTTPException(
            status_code=404,
            detail="officecli is not installed.",
        )

    workspace = await get_agent_for_request(request)
    coding_dir = await get_project_dir_for_request(request, workspace)
    target = _resolve_file_path_from_url(
        body.url,
        coding_dir,
        workspace.workspace_dir,
    )

    def _run_outline() -> dict:
        result = subprocess.run(  # noqa: S603
            [
                _officecli_bin(),
                "view",
                str(target),
                "outline",
                "--json",
            ],
            capture_output=True,
            text=True,
            timeout=_OFFICECLI_TIMEOUT,
            check=False,
        )
        if result.returncode != 0:
            raise HTTPException(
                status_code=500,
                detail=f"officecli: {result.stderr}",
            )
        try:
            return json.loads(result.stdout)
        except json.JSONDecodeError as exc:
            raise HTTPException(
                status_code=500,
                detail="officecli returned invalid JSON",
            ) from exc

    try:
        data = await asyncio.to_thread(_run_outline)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    return data


@router.post(
    "/office-issues",
    summary="Detect document issues (formatting/accessibility)",
)
async def office_issues(
    request: Request,
    body: ConvertOfficeRequest,
) -> dict:
    """Detect issues in an Office document (formatting, accessibility).

    Uses ``officecli view <file> issues --json``.
    """
    if not _is_officecli_available():
        raise HTTPException(
            status_code=404,
            detail="officecli is not installed.",
        )

    workspace = await get_agent_for_request(request)
    coding_dir = await get_project_dir_for_request(request, workspace)
    target = _resolve_file_path_from_url(
        body.url,
        coding_dir,
        workspace.workspace_dir,
    )

    def _run_issues() -> dict:
        result = subprocess.run(  # noqa: S603
            [
                _officecli_bin(),
                "view",
                str(target),
                "issues",
                "--json",
            ],
            capture_output=True,
            text=True,
            timeout=_OFFICECLI_TIMEOUT,
            check=False,
        )
        if result.returncode != 0:
            raise HTTPException(
                status_code=500,
                detail=f"officecli: {result.stderr}",
            )
        try:
            return json.loads(result.stdout)
        except json.JSONDecodeError as exc:
            raise HTTPException(
                status_code=500,
                detail="officecli returned invalid JSON",
            ) from exc

    try:
        data = await asyncio.to_thread(_run_issues)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    return data


@router.get(
    "/watch",
    summary="SSE stream for agent workspace file changes",
)
async def watch_workspace_files(
    request: Request,
    root: str = Query(default="project"),
) -> StreamingResponse:
    """Server-Sent Events that emit file-change notifications.

    Each SSE payload has the form::

        {"type": "file_change", "events": [{"change": "modified", "path": "..."}]}  # noqa: E501

    A heartbeat comment (``": heartbeat"``) is sent every 30 s when idle.
    """
    workspace = await get_agent_for_request(request)
    watch_dir = await _resolve_files_root(request, workspace, root)

    return StreamingResponse(
        workspace_watch_events(request, watch_dir),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


async def workspace_watch_events(
    request: Request,
    watch_dir: Path,
) -> AsyncIterator[str]:
    """Yield workspace file changes without cancelling the watcher on idle."""
    yield 'data: {"type": "connected"}\n\n'
    watcher = awatch(
        watch_dir,
        rust_timeout=_WATCH_POLL_TIMEOUT_MS,
        yield_on_timeout=True,
    )
    last_emit = asyncio.get_running_loop().time()
    try:
        while True:
            if await request.is_disconnected():
                break
            try:
                raw_changes = await watcher.__anext__()
            except (
                StopAsyncIteration,
                asyncio.CancelledError,
                GeneratorExit,
            ):
                break

            events = []
            for change_type, path in raw_changes:
                try:
                    rel = Path(path).relative_to(watch_dir)
                except ValueError:
                    continue
                if _should_skip(rel.parts):
                    continue
                change_name = (
                    "added"
                    if change_type is Change.added
                    else (
                        "deleted"
                        if change_type is Change.deleted
                        else "modified"
                    )
                )
                events.append(
                    {"change": change_name, "path": rel.as_posix()},
                )

            now = asyncio.get_running_loop().time()
            if events:
                payload = json.dumps(
                    {"type": "file_change", "events": events},
                    ensure_ascii=False,
                )
                yield f"data: {payload}\n\n"
                last_emit = now
            elif now - last_emit >= _WATCH_HEARTBEAT_SECONDS:
                yield ": heartbeat\n\n"
                last_emit = now
    except (asyncio.CancelledError, GeneratorExit):
        pass
    finally:
        try:
            await watcher.aclose()
        except Exception:
            pass


@router.get(
    "/memory",
    response_model=list[MdFileInfo],
    summary="List memory files",
    description="List all memory files (uses active agent)",
)
async def list_memory_files(
    request: Request,
    section: Literal["daily", "digest"] | None = Query(default=None),
) -> list[MdFileInfo]:
    """List memory directory markdown files."""
    try:
        workspace = await get_agent_for_request(request)
        workspace_manager = AgentMdManager(
            str(workspace.workspace_dir),
            agent_id=workspace.agent_id,
        )
        raw_files = await asyncio.to_thread(
            workspace_manager.list_memory_mds,
            section,
        )
        files = [MdFileInfo.model_validate(file) for file in raw_files]
        return files
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.get(
    "/memory/{md_path:path}",
    response_model=MdFileContent,
    summary="Read a memory file",
    description="Read a memory markdown file (uses active agent)",
)
async def read_memory_file(
    md_path: str,
    request: Request,
    section: Literal["daily", "digest"] | None = Query(default=None),
) -> MdFileContent:
    """Read a memory directory markdown file."""
    try:
        workspace = await get_agent_for_request(request)
        workspace_manager = AgentMdManager(
            str(workspace.workspace_dir),
            agent_id=workspace.agent_id,
        )
        content = await asyncio.to_thread(
            workspace_manager.read_memory_md,
            md_path,
            section,
        )
        return MdFileContent(content=content)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.put(
    "/memory/{md_path:path}",
    response_model=dict,
    summary="Write a memory file",
    description="Create or update a memory file (uses active agent)",
)
async def write_memory_file(
    md_path: str,
    body: MdFileContent,
    request: Request,
    section: Literal["daily", "digest"] | None = Query(default=None),
) -> dict:
    """Write a memory directory markdown file."""
    try:
        workspace = await get_agent_for_request(request)
        workspace_manager = AgentMdManager(
            str(workspace.workspace_dir),
            agent_id=workspace.agent_id,
        )
        await asyncio.to_thread(
            workspace_manager.write_memory_md,
            md_path,
            body.content,
            section,
        )
        return {"written": True}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.get(
    "/language",
    summary="Get agent language",
    description="Get the language setting for agent MD files.",
)
async def get_agent_language(request: Request) -> dict:
    """Get agent language setting for current agent."""
    workspace = await get_agent_for_request(request)
    agent_config = load_agent_config(workspace.agent_id)
    return {
        "language": agent_config.language,
        "agent_id": workspace.agent_id,
    }


@router.put(
    "/language",
    summary="Update agent language",
    description=(
        "Update the language for agent MD files. "
        "Optionally copies MD files for the new language to agent workspace."
    ),
)
async def put_agent_language(
    request: Request,
    body: dict = Body(
        ...,
        description='Language setting, e.g. {"language": "id"}',
    ),
) -> dict:
    """
    Update agent language and optionally re-copy MD files to agent workspace.
    """
    language = (body.get("language") or "").strip().lower()
    valid = SUPPORTED_AGENT_LANGUAGES
    if language not in valid:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Invalid language '{language}'. "
                f"Must be one of: {', '.join(sorted(valid))}"
            ),
        )

    workspace = await get_agent_for_request(request)
    agent_id = workspace.agent_id

    agent_config = load_agent_config(agent_id)
    old_language = agent_config.language

    agent_config.language = language
    save_agent_config(agent_id, agent_config)

    copied_files: list[str] = []
    if old_language != language:
        copied_files = copy_workspace_md_files(
            language,
            workspace.workspace_dir,
            md_template_id=get_workspace_md_template_id(
                agent_config.template_id
                or ("qa" if agent_id == BUILTIN_QA_AGENT_ID else None),
            ),
            only_if_missing=False,
        )

    return {
        "language": language,
        "copied_files": copied_files,
        "agent_id": agent_id,
    }


@router.get(
    "/audio-mode",
    summary="Get audio mode",
    description=(
        "Get the audio handling mode for incoming voice messages. "
        'Values: "auto", "native".'
    ),
)
async def get_audio_mode() -> dict:
    """Get audio mode setting."""
    config = load_config()
    return {"audio_mode": config.agents.audio_mode}


@router.put(
    "/audio-mode",
    summary="Update audio mode",
    description=(
        "Update how incoming audio/voice messages are handled. "
        '"auto": transcribe if provider available, else file placeholder; '
        '"native": send audio directly to model (may need ffmpeg).'
    ),
)
async def put_audio_mode(
    body: dict = Body(
        ...,
        description='Audio mode, e.g. {"audio_mode": "auto"}',
    ),
) -> dict:
    """Update audio mode setting."""
    raw = body.get("audio_mode")
    audio_mode = (str(raw) if raw is not None else "").strip().lower()
    valid = {"auto", "native"}
    if audio_mode not in valid:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Invalid audio_mode '{audio_mode}'. "
                f"Must be one of: {', '.join(sorted(valid))}"
            ),
        )

    def apply_audio_mode(config: Any) -> None:
        config.agents.audio_mode = audio_mode

    await run_sync_io(mutate_config, apply_audio_mode)
    return {"audio_mode": audio_mode}


@router.get(
    "/transcription-provider-type",
    summary="Get transcription provider type",
    description=(
        "Get the transcription provider type. "
        'Values: "disabled", "whisper_api", "local_whisper".'
    ),
)
async def get_transcription_provider_type() -> dict:
    """Get transcription provider type setting."""
    config = load_config()
    return {
        "transcription_provider_type": (
            config.agents.transcription_provider_type
        ),
    }


@router.put(
    "/transcription-provider-type",
    summary="Set transcription provider type",
    description=(
        "Set the transcription provider type. "
        '"disabled": no transcription; '
        '"whisper_api": remote Whisper endpoint; '
        '"local_whisper": locally installed openai-whisper.'
    ),
)
async def put_transcription_provider_type(
    body: dict = Body(
        ...,
        description=(
            "Provider type, e.g. "
            '{"transcription_provider_type": "whisper_api"}'
        ),
    ),
) -> dict:
    """Set the transcription provider type."""
    raw = body.get("transcription_provider_type")
    provider_type = (str(raw) if raw is not None else "").strip().lower()
    valid = {"disabled", "whisper_api", "local_whisper"}
    if provider_type not in valid:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Invalid transcription_provider_type '{provider_type}'. "
                f"Must be one of: {', '.join(sorted(valid))}"
            ),
        )

    def apply_provider_type(config: Any) -> None:
        config.agents.transcription_provider_type = provider_type

    await run_sync_io(mutate_config, apply_provider_type)
    return {"transcription_provider_type": provider_type}


@router.get(
    "/local-whisper-status",
    summary="Check local whisper availability",
    description=(
        "Check whether the local whisper provider can be used. "
        "Returns availability of ffmpeg and openai-whisper."
    ),
)
async def get_local_whisper_status() -> dict:
    """Check local whisper dependencies."""
    from ...agents.utils.audio_transcription import (
        check_local_whisper_available,
    )

    return check_local_whisper_available()


@router.get(
    "/transcription-providers",
    summary="List transcription providers",
    description=(
        "List providers capable of audio transcription (Whisper API). "
        "Returns available providers and the configured selection."
    ),
)
async def get_transcription_providers() -> dict:
    """List transcription-capable providers and configured selection."""
    from ...agents.utils.audio_transcription import (
        get_configured_transcription_provider_id,
        list_transcription_providers,
    )

    return {
        "providers": list_transcription_providers(),
        "configured_provider_id": (get_configured_transcription_provider_id()),
    }


@router.put(
    "/transcription-provider",
    summary="Set transcription provider",
    description=(
        "Set the provider to use for audio transcription. "
        'Use empty string "" to unset.'
    ),
)
async def put_transcription_provider(
    body: dict = Body(
        ...,
        description=(
            'Provider ID, e.g. {"provider_id": "openai"} '
            'or {"provider_id": ""} to unset'
        ),
    ),
) -> dict:
    """Set the transcription provider."""
    provider_id = (body.get("provider_id") or "").strip()

    def apply_provider(config: Any) -> None:
        config.agents.transcription_provider_id = provider_id

    await run_sync_io(mutate_config, apply_provider)
    return {"provider_id": provider_id}


@router.post(
    "/transcribe",
    summary="Transcribe audio to text",
    description=(
        "Transcribe an uploaded audio file "
        "using the configured Whisper provider. "
        "Returns the transcribed text."
    ),
)
async def post_transcribe_audio(
    file: UploadFile = File(..., description="Audio file to transcribe"),
) -> dict:
    """Transcribe uploaded audio file using configured Whisper provider."""
    from ...agents.utils.audio_transcription import transcribe_audio

    # Check transcription is enabled
    config = load_config()
    provider_type = config.agents.transcription_provider_type
    if provider_type == "disabled":
        raise HTTPException(
            status_code=400,
            detail={
                "code": "TRANSCRIPTION_DISABLED",
                "message": (
                    "Transcription is disabled. "
                    "Configure a transcription provider in Settings."
                ),
            },
        )

    # Validate file type
    allowed_extensions = {
        ".webm",
        ".mp4",
        ".m4a",
        ".wav",
        ".mp3",
        ".ogg",
        ".flac",
    }
    suffix = (
        os.path.splitext(file.filename or "audio.webm")[1].lower() or ".webm"
    )
    if suffix not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail={
                "code": "UNSUPPORTED_FILE_TYPE",
                "message": (
                    f"Unsupported file type: {suffix}. "
                    f"Allowed: {', '.join(sorted(allowed_extensions))}"
                ),
            },
        )

    data = await file.read()
    check_upload_size(data)

    # Save uploaded file to temp directory
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        text = await transcribe_audio(tmp_path)
        if text is None:
            raise HTTPException(
                status_code=500,
                detail="Transcription failed. Check provider configuration.",
            )
        return {"text": text}
    finally:
        # Clean up temp file
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


@router.post(
    "/embedding/test",
    response_model=EmbeddingTestResponse,
    summary="Test embedding configuration",
    description=(
        "Create an AgentScope embedding model, perform a real request, and "
        "validate the returned dimensions"
    ),
)
async def test_embedding_configuration(
    embedding_config: EmbeddingModelConfig = Body(...),
    request: Request = None,
) -> EmbeddingTestResponse:
    """Test unsaved embedding settings and stage the model for hot apply."""
    workspace = await get_agent_for_request(request)
    memory_manager = workspace.memory_manager
    if memory_manager is not None and hasattr(
        memory_manager,
        "test_and_stage_embedding",
    ):
        result = await memory_manager.test_and_stage_embedding(
            embedding_config,
        )
    else:
        _model, result = await test_embedding_model(embedding_config)

    message = result.message
    if embedding_config.api_key:
        message = message.replace(embedding_config.api_key, "***")
    return EmbeddingTestResponse(
        success=result.success,
        configured_dimensions=result.configured_dimensions,
        actual_dimensions=result.actual_dimensions,
        latency_ms=result.latency_ms,
        message=message,
    )


@router.get(
    "/running-config",
    response_model=AgentsRunningConfig,
    summary="Get agent running config",
    description="Get running configuration for active agent",
)
async def get_agents_running_config(
    request: Request,
) -> AgentsRunningConfig:
    """Get agent running configuration."""
    workspace = await get_agent_for_request(request)
    agent_config = await run_sync_io(load_agent_config, workspace.agent_id)
    running = agent_config.running or AgentsRunningConfig()
    running.approval_level = getattr(agent_config, "approval_level", "AUTO")
    return running


class _ConfigRollbackConflict(RuntimeError):
    """Raised when a field changed again after this request persisted it."""

    def __init__(self, paths: list[str]):
        super().__init__("configuration changed concurrently")
        self.paths = paths


def _conditionally_restore_config_changes(
    current: BaseModel,
    before: BaseModel,
    submitted: BaseModel,
) -> None:
    """Three-way rollback without overwriting unrelated concurrent edits."""
    candidate = current.model_copy(deep=True)
    conflicts: list[str] = []

    def restore(
        target: BaseModel,
        old: BaseModel,
        saved: BaseModel,
        prefix: str,
    ) -> None:
        for name in type(saved).model_fields:
            old_value = getattr(old, name)
            saved_value = getattr(saved, name)
            if old_value == saved_value:
                continue
            current_value = getattr(target, name)
            path = f"{prefix}.{name}" if prefix else name
            if (
                isinstance(current_value, BaseModel)
                and isinstance(old_value, BaseModel)
                and isinstance(saved_value, BaseModel)
                and type(current_value) is type(old_value) is type(saved_value)
            ):
                restore(current_value, old_value, saved_value, path)
            elif current_value == saved_value:
                setattr(target, name, copy.deepcopy(old_value))
            else:
                conflicts.append(path)

    restore(candidate, before, submitted, "")
    if conflicts:
        raise _ConfigRollbackConflict(conflicts)
    for field_name in type(current).model_fields:
        setattr(current, field_name, getattr(candidate, field_name))


async def _apply_embedding_runtime(
    memory_manager: Any,
    embedding_config: EmbeddingModelConfig,
    agent_id: str,
) -> bool:
    """Apply an embedding config to a running memory manager."""
    if hasattr(memory_manager, "apply_tested_embedding"):
        try:
            if await memory_manager.apply_tested_embedding(embedding_config):
                return True
        except Exception as exc:
            logger.warning(
                "Embedding hot update failed for agent '%s': %s",
                agent_id,
                exc,
                exc_info=True,
            )
    if hasattr(memory_manager, "reload_embedding_config"):
        try:
            return bool(await memory_manager.reload_embedding_config())
        except Exception as exc:
            logger.warning(
                "Embedding runtime reload failed for agent '%s': %s",
                agent_id,
                exc,
                exc_info=True,
            )
    return False


async def _rollback_embedding_update(
    agent_id: str,
    memory_manager: Any,
    before: BaseModel,
    submitted: BaseModel,
) -> None:
    """Roll back persistence and runtime after an embedding update fails."""
    rollback_conflict: _ConfigRollbackConflict | None = None

    def rollback_config(current_config: BaseModel) -> None:
        _conditionally_restore_config_changes(
            current_config,
            before,
            submitted,
        )

    try:
        await update_agent_config_async(agent_id, rollback_config)
    except _ConfigRollbackConflict as exc:
        rollback_conflict = exc

    runtime_restored = False
    if hasattr(memory_manager, "reload_embedding_config"):
        try:
            runtime_restored = bool(
                await memory_manager.reload_embedding_config(),
            )
        except Exception:
            logger.exception(
                "Failed to restore the previous embedding runtime "
                "for agent '%s'",
                agent_id,
            )

    raise HTTPException(
        status_code=409 if rollback_conflict else 503,
        detail={
            "message": (
                "Embedding configuration was not applied; "
                + (
                    "rollback was skipped because the configuration "
                    "changed concurrently"
                    if rollback_conflict
                    else "the persisted changes were rolled back"
                )
            ),
            "persisted": rollback_conflict is not None,
            "runtime_applied": False,
            "runtime_restored": runtime_restored,
            "conflicts": rollback_conflict.paths if rollback_conflict else [],
        },
    )


@router.put(
    "/running-config",
    response_model=AgentsRunningConfig,
    summary="Update agent running config",
    description="Update running configuration for active agent",
)
async def put_agents_running_config(
    running_config: AgentsRunningConfig = Body(
        ...,
        description="Updated agent running configuration",
    ),
    request: Request = None,
) -> AgentsRunningConfig:
    """Update agent running configuration."""
    workspace = await get_agent_for_request(request)
    memory_manager = workspace.memory_manager
    workspace_dir = getattr(workspace, "workspace_dir", ".")
    config_path = Path(workspace_dir) / "agent.json"
    async with get_path_lock(config_path):
        old_agent_config = None
        embedding_changed = False
        memory_manager_backend_changed = False
        new_embedding_config = (
            running_config.reme_light_memory_config.embedding_model_config
        )
        new_memory_manager_backend = running_config.memory_manager_backend

        def persist_running_config(agent_config):
            nonlocal old_agent_config, embedding_changed
            nonlocal memory_manager_backend_changed
            old_agent_config = agent_config.model_copy(deep=True)
            old_running_config = agent_config.running or AgentsRunningConfig()
            memory_manager_backend_changed = (
                old_running_config.memory_manager_backend
                != new_memory_manager_backend
            )
            old_memory_config = old_running_config.reme_light_memory_config
            old_embedding_config = old_memory_config.embedding_model_config
            vector_space_changed = embedding_vector_space_fingerprint(
                old_embedding_config,
            ) != embedding_vector_space_fingerprint(new_embedding_config)
            running_config.reme_light_memory_config.needs_reindex = (
                old_memory_config.needs_reindex or vector_space_changed
            )
            embedding_changed = old_embedding_config != new_embedding_config
            if (
                embedding_changed
                and not memory_manager_backend_changed
                and new_memory_manager_backend == "remelight"
                and memory_manager is not None
                and getattr(memory_manager, "is_reindexing", False) is True
            ):
                raise HTTPException(
                    status_code=409,
                    detail=(
                        "Embedding configuration cannot change while the "
                        "memory index is rebuilding"
                    ),
                )
            if running_config.approval_level is not None:
                agent_config.approval_level = running_config.approval_level
            running_config.approval_level = None
            agent_config.running = running_config

        agent_config = await update_agent_config_async(
            workspace.agent_id,
            persist_running_config,
        )

        if (
            embedding_changed
            and not memory_manager_backend_changed
            and new_memory_manager_backend == "remelight"
            and memory_manager is not None
        ):
            embedding_updated = await _apply_embedding_runtime(
                memory_manager,
                new_embedding_config,
                workspace.agent_id,
            )
            if not embedding_updated:
                assert old_agent_config is not None
                await _rollback_embedding_update(
                    workspace.agent_id,
                    memory_manager,
                    old_agent_config,
                    agent_config,
                )

    schedule_agent_reload(request, workspace.agent_id)

    running_config.approval_level = agent_config.approval_level
    return running_config


@router.get(
    "/system-prompt-files",
    response_model=list[str],
    summary="Get system prompt files",
    description="Get system prompt files for active agent",
)
async def get_system_prompt_files(
    request: Request,
) -> list[str]:
    """Get list of enabled system prompt files."""
    workspace = await get_agent_for_request(request)
    agent_config = load_agent_config(workspace.agent_id)
    return agent_config.system_prompt_files or []


@router.put(
    "/system-prompt-files",
    response_model=list[str],
    summary="Update system prompt files",
    description="Update system prompt files for active agent",
)
async def put_system_prompt_files(
    files: list[str] = Body(
        ...,
        description="Markdown filenames to load into system prompt",
    ),
    request: Request = None,
) -> list[str]:
    """Update list of enabled system prompt files."""
    workspace = await get_agent_for_request(request)
    agent_config = load_agent_config(workspace.agent_id)
    agent_config.system_prompt_files = files
    save_agent_config(workspace.agent_id, agent_config)

    schedule_agent_reload(request, workspace.agent_id)

    return files


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _validate_zip_data(data: bytes, workspace_dir: Path) -> None:
    """Ensure *data* is a valid zip without path-traversal entries."""
    if not zipfile.is_zipfile(io.BytesIO(data)):
        raise HTTPException(
            status_code=400,
            detail="Uploaded file is not a valid zip archive",
        )
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        for name in zf.namelist():
            resolved = (workspace_dir / name).resolve()
            if not str(resolved).startswith(str(workspace_dir)):
                raise HTTPException(
                    status_code=400,
                    detail=f"Zip contains unsafe path: {name}",
                )


def _extract_and_merge_zip(data: bytes, workspace_dir: Path) -> None:
    """Extract zip data and merge into workspace_dir (blocking operation)."""
    tmp_dir = None
    try:
        tmp_dir = Path(tempfile.mkdtemp(prefix="qwenpaw_upload_"))
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            zf.extractall(tmp_dir)

        top_entries = list(tmp_dir.iterdir())
        extract_root = tmp_dir
        if len(top_entries) == 1 and top_entries[0].is_dir():
            extract_root = top_entries[0]

        workspace_dir.mkdir(parents=True, exist_ok=True)

        for item in extract_root.iterdir():
            dest = workspace_dir / item.name
            if item.is_file():
                shutil.copy2(item, dest)
            else:
                if dest.exists() and dest.is_file():
                    dest.unlink()
                shutil.copytree(item, dest, dirs_exist_ok=True)
    finally:
        if tmp_dir and tmp_dir.is_dir():
            shutil.rmtree(tmp_dir, ignore_errors=True)


def _validate_and_extract_zip(data: bytes, workspace_dir: Path) -> None:
    """Validate and extract zip data (blocking operation)."""
    _validate_zip_data(data, workspace_dir)
    _extract_and_merge_zip(data, workspace_dir)


# ---------------------------------------------------------------------------
# Workspace Download/Upload Endpoints
# ---------------------------------------------------------------------------


@router.get(
    "/download",
    summary="Download workspace as zip",
    description=(
        "Package the entire agent workspace into a zip archive and stream "
        "it back as a downloadable file."
    ),
    responses={
        200: {
            "content": {"application/zip": {}},
            "description": "Zip archive of agent workspace",
        },
    },
)
async def download_workspace(request: Request):
    """Stream agent workspace as a zip file."""

    agent = await get_agent_for_request(request)
    workspace_dir = agent.workspace_dir

    if not workspace_dir.is_dir():
        raise HTTPException(
            status_code=404,
            detail=f"Workspace does not exist: {workspace_dir}",
        )

    buf = await asyncio.to_thread(_zip_directory, workspace_dir)

    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    filename = f"qwenpaw_workspace_{agent.agent_id}_{timestamp}.zip"

    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
        },
    )


@router.post(
    "/upload",
    response_model=dict,
    summary="Upload zip and merge into workspace",
    description=(
        "Upload a zip archive.  Paths present in the zip are merged into "
        "agent workspace (files overwritten, dirs merged).  Paths not in "
        "the zip are left unchanged (e.g. qwenpaw.db, runtime dirs). "
        "Download packs the entire workspace; upload only "
        "overwrites/merges zip contents."
    ),
)
async def upload_workspace(
    request: Request,
    file: UploadFile = File(
        ...,
        description="Zip archive to merge into agent workspace",
    ),
) -> dict:
    """
    Merge uploaded zip contents into agent workspace (overwrite, not clear).
    """

    if file.content_type and file.content_type not in (
        "application/zip",
        "application/x-zip-compressed",
        "application/octet-stream",
    ):
        raise HTTPException(
            status_code=400,
            detail=(
                f"Expected a zip file, got content-type: {file.content_type}"
            ),
        )

    agent = await get_agent_for_request(request)
    workspace_dir = agent.workspace_dir
    data = await file.read()

    try:
        await asyncio.to_thread(_validate_and_extract_zip, data, workspace_dir)
        return {"success": True}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to merge workspace: {exc}",
        ) from exc


@router.get("/commands/available")
async def get_available_commands(request: Request):
    """Return all slash commands registered for the workspace.

    Merges built-in system commands with plugin-registered ones
    so the frontend can dynamically populate the slash menu.
    """
    agent = await get_agent_for_request(request)
    registry = getattr(
        getattr(agent, "plugins", None),
        "slash_command_registry",
        None,
    )
    commands = []
    if registry is not None:
        for name in registry.names():
            match = registry.resolve(f"/{name}")
            desc = ""
            category = ""
            if match:
                spec, _ = match
                desc = spec.help_text or ""
                category = spec.category or ""
            commands.append(
                {
                    "name": name,
                    "description": desc,
                    "category": category,
                },
            )
    return ORJSONResponse({"commands": commands})
